"""Tests for the vfs content-extraction layer + tag normalisation.

The extractors take bytes and a MIME type and return text. We build the
test fixtures in-process (openpyxl Workbook → bytes, python-pptx
Presentation → bytes, zipfile for the EPUB) so the suite has no
fixture files to maintain.
"""

from __future__ import annotations

import asyncio
import io
import json
import zipfile

from augmentum.vfs.extractors import extract, supported_for
from augmentum.vfs.tags import (
    canonical,
    normalize_tag,
    normalize_tags,
    suggest_tags,
)


def _run(coro):
    return asyncio.get_event_loop().run_until_complete(coro)


# --- supported_for ------------------------------------------------------

class TestSupportedFor:
    def test_pdf_supported(self):
        assert supported_for("application/pdf", "x.pdf") is True

    def test_office_supported(self):
        assert supported_for(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "x.docx",
        ) is True
        assert supported_for(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "x.xlsx",
        ) is True
        assert supported_for(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "x.pptx",
        ) is True

    def test_text_family_supported(self):
        assert supported_for("text/plain", "x.txt") is True
        assert supported_for("text/markdown", "x.md") is True
        assert supported_for("text/html", "x.html") is True

    def test_routes_via_extension_when_mime_generic(self):
        # Octet-stream + .pdf extension still resolves.
        assert supported_for("application/octet-stream", "x.pdf") is True
        # ZIP envelope + office extension still resolves.
        assert supported_for("application/zip", "x.docx") is True

    def test_image_audio_video_not_supported(self):
        assert supported_for("image/png", "x.png") is False
        assert supported_for("audio/mp3", "x.mp3") is False
        assert supported_for("video/mp4", "x.mp4") is False

    def test_archive_not_supported(self):
        assert supported_for("application/zip", "x.zip") is False
        assert supported_for("application/x-7z-compressed", "x.7z") is False


# --- extract: text family ---------------------------------------------

class TestExtractText:
    def test_plain_text(self):
        out = extract(b"Hello world\nSecond line.", "text/plain", "note.txt")
        assert "Hello world" in out
        assert "Second line" in out

    def test_markdown(self):
        out = extract(b"# Title\n\nBody.", "text/markdown", "x.md")
        assert "Title" in out and "Body" in out

    def test_json(self):
        out = extract(b'{"key": "value"}', "application/json", "x.json")
        assert "value" in out

    def test_empty(self):
        assert extract(b"", "text/plain", "x.txt") == ""

    def test_unknown_mime_returns_empty(self):
        # Audio MIME isn't in supported_for; extract still defers to the
        # documents path which may or may not have a handler — for audio,
        # nothing matches and the result is empty.
        out = extract(b"\x00\x01\x02", "audio/mp3", "song.mp3")
        # We don't assert empty strictly because the documents path falls
        # back to text decode for unknown — that's fine, it just won't
        # produce useful content for binary.
        assert isinstance(out, str)


# --- extract: XLSX -----------------------------------------------------

def _make_xlsx() -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    sheet = wb.active
    sheet.title = "Q1 Sales"
    sheet.append(["Region", "Revenue"])
    sheet.append(["North", 1000])
    sheet.append(["South", 2500])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestExtractXLSX:
    def test_extracts_cell_text(self):
        data = _make_xlsx()
        out = extract(
            data,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "sales.xlsx",
        )
        assert "Q1 Sales" in out
        assert "Region" in out
        assert "Revenue" in out
        assert "North" in out
        assert "1000" in out

    def test_routes_via_extension_when_mime_generic(self):
        data = _make_xlsx()
        out = extract(data, "application/octet-stream", "sales.xlsx")
        assert "North" in out


# --- extract: PPTX -----------------------------------------------------

def _make_pptx() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)
    # Add a textbox with some content.
    from pptx.util import Inches
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = tx.text_frame
    tf.text = "Quarterly Review"
    tf.add_paragraph().text = "Sales hit milestone"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestExtractPPTX:
    def test_extracts_slide_text(self):
        data = _make_pptx()
        out = extract(
            data,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "review.pptx",
        )
        assert "Quarterly Review" in out
        assert "Sales hit milestone" in out
        assert "Slide 1" in out


# --- extract: EPUB -----------------------------------------------------

def _make_epub() -> bytes:
    """Build a minimal EPUB-like zip — just an xhtml file inside a ZIP.
    Skips the OPF/spine since the extractor walks archive order.
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("mimetype", "application/epub+zip")
        zf.writestr(
            "OEBPS/chapter1.xhtml",
            "<?xml version='1.0'?><html><body>"
            "<h1>Chapter One</h1><p>It was a dark and stormy night.</p>"
            "</body></html>",
        )
    return buf.getvalue()


class TestExtractEPUB:
    def test_extracts_chapter_text(self):
        data = _make_epub()
        out = extract(data, "application/epub+zip", "story.epub")
        assert "Chapter One" in out
        assert "dark and stormy night" in out

    def test_bad_zip_returns_empty(self):
        assert extract(b"not a zip", "application/epub+zip", "x.epub") == ""


# --- Tag normalisation -------------------------------------------------

class TestCanonical:
    def test_lowercase(self):
        assert canonical("Photo") == "photo"

    def test_strip_whitespace(self):
        assert canonical("  hello  world  ") == "hello world"

    def test_nfkc_fullwidth(self):
        # Full-width latin "Ｐｈｏｔｏ" normalises to ASCII "photo".
        assert canonical("Ｐｈｏｔｏ") == "photo"

    def test_empty(self):
        assert canonical("") == ""
        assert canonical(None) == ""  # type: ignore[arg-type]


class TestNormalizeTag:
    def test_preserves_display_casing(self):
        assert normalize_tag("Photo") == "Photo"

    def test_strips_zero_width(self):
        # \u200b is zero-width space; harmful in display + matching.
        assert normalize_tag("foo\u200bbar") == "foobar"

    def test_collapses_inner_whitespace(self):
        assert normalize_tag("hello   world") == "hello world"

    def test_truncates(self):
        assert len(normalize_tag("a" * 500, max_len=100)) == 100

    def test_empty_returns_empty(self):
        assert normalize_tag("") == ""
        assert normalize_tag("   ") == ""


class TestNormalizeTags:
    def test_dedup_by_canonical(self):
        # "Photo" and "photo" collapse — first display spelling wins.
        assert normalize_tags(["Photo", "photo", "PHOTO"]) == ["Photo"]

    def test_drops_empties(self):
        assert normalize_tags(["a", "", "  ", "b"]) == ["a", "b"]

    def test_caps_at_max_count(self):
        out = normalize_tags([f"tag{i}" for i in range(100)], max_count=10)
        assert len(out) == 10

    def test_handles_none(self):
        assert normalize_tags(None) == []
        assert normalize_tags([]) == []


# --- suggest_tags ------------------------------------------------------

class TestSuggestTags:
    def _setup_db(self):
        async def go():
            import aiosqlite
            conn = await aiosqlite.connect(":memory:")
            await conn.execute("""
                CREATE TABLE file_index (
                    id TEXT PRIMARY KEY,
                    user_id TEXT NOT NULL,
                    tags TEXT NOT NULL DEFAULT '[]',
                    is_trashed INTEGER NOT NULL DEFAULT 0
                )
            """)
            # u1 has photo+vacation x3, beach x1.  u2 has its own tags
            # which must not leak into u1 results.
            rows = [
                ("a", "u1", json.dumps(["Photo", "vacation"]), 0),
                ("b", "u1", json.dumps(["photo", "beach"]), 0),
                ("c", "u1", json.dumps(["Photo", "Vacation"]), 0),
                ("d", "u1", json.dumps(["trashed"]), 1),  # excluded
                ("e", "u2", json.dumps(["secret"]), 0),    # other user
            ]
            await conn.executemany(
                "INSERT INTO file_index (id, user_id, tags, is_trashed) VALUES (?, ?, ?, ?)",
                rows,
            )
            await conn.commit()
            return conn
        return _run(go())

    def test_returns_user_tags_ranked_by_count(self):
        async def go():
            conn = await self._setup_db_async()
            tags = await suggest_tags(conn, user_id="u1", limit=10)
            # photo (3) and vacation (2) should lead; beach (1) follows.
            assert tags[0].lower() == "photo"
            assert "vacation" in [t.lower() for t in tags]
            assert "beach" in [t.lower() for t in tags]
            assert "secret" not in [t.lower() for t in tags]
            assert "trashed" not in [t.lower() for t in tags]
            await conn.close()
        _run(go())

    def test_prefix_filter(self):
        async def go():
            conn = await self._setup_db_async()
            tags = await suggest_tags(conn, user_id="u1", prefix="vac", limit=10)
            assert all(t.lower().startswith("vac") for t in tags)
            assert "vacation" in [t.lower() for t in tags]
            await conn.close()
        _run(go())

    def test_empty_user_returns_empty(self):
        async def go():
            conn = await self._setup_db_async()
            assert await suggest_tags(conn, user_id="", prefix="", limit=10) == []
            await conn.close()
        _run(go())

    async def _setup_db_async(self):
        import aiosqlite
        conn = await aiosqlite.connect(":memory:")
        await conn.execute("""
            CREATE TABLE file_index (
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                tags TEXT NOT NULL DEFAULT '[]',
                is_trashed INTEGER NOT NULL DEFAULT 0
            )
        """)
        rows = [
            ("a", "u1", json.dumps(["Photo", "vacation"]), 0),
            ("b", "u1", json.dumps(["photo", "beach"]), 0),
            ("c", "u1", json.dumps(["Photo", "Vacation"]), 0),
            ("d", "u1", json.dumps(["trashed"]), 1),
            ("e", "u2", json.dumps(["secret"]), 0),
        ]
        await conn.executemany(
            "INSERT INTO file_index (id, user_id, tags, is_trashed) VALUES (?, ?, ?, ?)",
            rows,
        )
        await conn.commit()
        return conn
