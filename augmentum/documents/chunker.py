"""Text extraction and chunking for document RAG.

Quality features:
- Section-aware chunking: splits on headings/paragraphs first, then by size
- Header enrichment: prepends filename + section context to each chunk
- PDF table extraction: pdfplumber tables → markdown format
- HTML via trafilatura: proper content extraction, not raw tag soup
- Parent-child chunks: small chunks for search precision, large for context
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass

from augmentum.utils.logging import get_logger

log = get_logger(__name__)

DEFAULT_CHUNK_SIZE = 1500       # chars (~375 tokens) — search chunk
DEFAULT_PARENT_CHUNK_SIZE = 4000  # chars (~1000 tokens) — context chunk
DEFAULT_CHUNK_OVERLAP = 200

# Heading patterns for section detection
_HEADING_RE = re.compile(
    r"^(?:"
    r"#{1,6}\s+.+|"              # Markdown headings
    r"[A-Z][A-Za-z0-9 :,\-]{2,80}$|"  # Title-case lines (standalone)
    r"\d+\.[\d.]*\s+[A-Z].{2,80}$"    # Numbered sections: "1.2 Overview"
    r")",
    re.MULTILINE,
)


@dataclass
class Chunk:
    """A single chunk of extracted text."""

    index: int
    text: str                          # Raw chunk content
    enriched_text: str = ""            # Header-enriched text (for embedding)
    page_num: int | None = None
    char_offset: int = 0
    section: str = ""                  # Section heading this chunk belongs to
    parent_index: int | None = None    # Index of parent (large) chunk


@dataclass
class _Section:
    """A structural section detected in the document."""
    heading: str
    content: str
    page_num: int | None = None


# =========================================================================
# Text extraction
# =========================================================================

def extract_text(data: bytes, mime_type: str, filename: str = "") -> list[tuple[str, int | None]]:
    """Extract text from file bytes.

    Returns a list of (text, page_number) tuples.
    page_number is None for non-paged formats.
    """
    if mime_type == "application/pdf":
        return _extract_pdf(data)
    if mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx(data)
    if mime_type in ("text/html", "application/xhtml+xml") or filename.endswith((".html", ".htm")):
        return _extract_html(data)
    # Office OOXML formats are ZIP envelopes — dispatch by filename FIRST, since
    # the upload path's mime is unreliable (a .pptx often arrives as
    # application/octet-stream or gets coerced to text/plain). Without this the
    # raw zip bytes fall through to the text fallback below and the model is fed
    # the zip's member-filename table (slide1.xml, theme1.xml, …) as "content".
    _lower = filename.lower()
    if _lower.endswith(".pptx") or mime_type == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        return _extract_pptx(data)
    if _lower.endswith(".xlsx") or mime_type == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        return _extract_xlsx(data)
    if mime_type.startswith("text/") or filename.endswith(
        (".md", ".markdown", ".txt", ".csv", ".json", ".log", ".py", ".js", ".ts",
         ".yaml", ".yml", ".toml", ".xml", ".rst"),
    ):
        text = data.decode("utf-8", errors="replace")
        return [(text, None)]
    # Fallback: try as text — but ONLY if the bytes actually look like text.
    # Binary blobs (zip-based office docs, images, archives) must NOT be decoded
    # and shipped to the model as garbage; return empty so the caller surfaces a
    # clean "could not extract text" error instead.
    if _looks_binary(data):
        log.warning("extract_text_binary_skipped", mime_type=mime_type, filename=filename)
        return []
    try:
        text = data.decode("utf-8", errors="replace")
        return [(text, None)]
    except Exception:
        log.warning("extract_text_failed", mime_type=mime_type, filename=filename)
        return []


# ZIP local-file-header / empty-archive / spanned magics — the OOXML envelope.
_ZIP_MAGICS = (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")


def _looks_binary(data: bytes) -> bool:
    """Heuristic: does this byte stream look like binary rather than text?

    Used to gate the last-resort text fallback so we never decode a zip/image/
    archive into the prompt. A NUL byte or a known binary magic is decisive; a
    high non-printable ratio in the head catches the rest.
    """
    if not data:
        return False
    head = data[:1024]
    if head.startswith(_ZIP_MAGICS) or b"\x00" in head:
        return True
    printable = sum(1 for b in head if b in (9, 10, 13) or 32 <= b <= 126)
    return printable / len(head) < 0.85


def _extract_pptx(data: bytes) -> list[tuple[str, int | None]]:
    """Extract per-slide text from a PPTX, one (text, slide_number) tuple each."""
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python_pptx_not_installed")
        return []
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception:
        log.warning("pptx_load_failed", exc_info=True)
        return []

    pages: list[tuple[str, int | None]] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text or "" for run in para.runs).strip()
                if line:
                    lines.append(line)
        if lines:
            pages.append(("\n".join(lines), i))
    return pages


def _extract_xlsx(data: bytes) -> list[tuple[str, int | None]]:
    """Flatten each worksheet into a markdown-ish table; one tuple per sheet."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        log.warning("openpyxl_not_installed")
        return []
    try:
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception:
        log.warning("xlsx_load_failed", exc_info=True)
        return []

    pages: list[tuple[str, int | None]] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= 500:
                rows.append("…")
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            pages.append((f"# {sheet.title}\n" + "\n".join(rows), None))
    wb.close()
    return pages


def _extract_pdf(data: bytes) -> list[tuple[str, int | None]]:
    """Extract text + tables from PDF via pdfplumber, falling back to pypdf."""
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages: list[tuple[str, int | None]] = []
            for i, page in enumerate(pdf.pages):
                parts: list[str] = []

                # Extract tables as markdown
                tables = page.extract_tables() or []
                table_text = ""
                if tables:
                    table_text = _tables_to_markdown(tables)

                # Extract regular text
                text = page.extract_text() or ""

                if text.strip():
                    parts.append(text)
                if table_text:
                    parts.append(table_text)

                combined = "\n\n".join(parts)
                if combined.strip():
                    pages.append((combined, i + 1))

            if pages:
                return pages
    except ImportError:
        pass
    except Exception:
        log.debug("pdfplumber_failed_trying_pypdf", exc_info=True)

    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append((text, i + 1))
        return pages
    except ImportError:
        log.warning("no_pdf_library_available")
        return _extract_pdf_fallback(data)


def _extract_pdf_fallback(data: bytes) -> list[tuple[str, int | None]]:
    """Last-resort PDF extraction."""
    try:
        text = data.decode("utf-8", errors="replace")
        printable = sum(1 for c in text[:1000] if c.isprintable() or c.isspace())
        if printable / max(len(text[:1000]), 1) > 0.8:
            return [(text, None)]
    except (UnicodeDecodeError, ZeroDivisionError):
        # Decode with errors='replace' shouldn't raise but be defensive;
        # printable ratio uses max() so ZeroDivisionError shouldn't fire
        # either — keep narrow to surface anything else.
        pass
    return []


def _extract_docx(data: bytes) -> list[tuple[str, int | None]]:
    """Extract text from DOCX preserving heading structure."""
    try:
        from docx import Document
    except ImportError:
        log.warning("python_docx_not_installed")
        return []

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Preserve heading markers for section detection
        style = (para.style.name or "").lower() if para.style else ""
        if "heading" in style:
            level = 1
            for ch in style:
                if ch.isdigit():
                    level = int(ch)
                    break
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    if not parts:
        return []
    return [("\n\n".join(parts), None)]


def _extract_html(data: bytes) -> list[tuple[str, int | None]]:
    """Extract clean text from HTML using trafilatura (already in deps)."""
    try:
        import trafilatura

        text = trafilatura.extract(
            data.decode("utf-8", errors="replace"),
            include_tables=True,
            include_links=False,
            include_images=False,
            output_format="txt",
        )
        if text and text.strip():
            return [(text, None)]
    except ImportError:
        pass
    except Exception:
        log.debug("trafilatura_html_extract_failed", exc_info=True)

    # Fallback: strip tags with regex
    html_str = data.decode("utf-8", errors="replace")
    text = re.sub(r"<[^>]+>", " ", html_str)
    text = re.sub(r"\s+", " ", text).strip()
    return [(text, None)] if text else []


def _tables_to_markdown(tables: list[list[list]]) -> str:
    """Convert pdfplumber table data to markdown tables."""
    parts: list[str] = []
    for table in tables:
        if not table or not table[0]:
            continue
        rows: list[str] = []
        for i, row in enumerate(table):
            cells = [str(c or "").replace("|", "\\|").strip() for c in row]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join("---" for _ in cells) + " |")
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


# =========================================================================
# Section detection
# =========================================================================

def _detect_sections(text: str, page_num: int | None = None) -> list[_Section]:
    """Split text into sections based on heading detection.

    Returns at least one section. If no headings are found, the entire
    text is returned as a single section with heading="".
    """
    lines = text.split("\n")
    sections: list[_Section] = []
    current_heading = ""
    current_lines: list[str] = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            current_lines.append("")
            continue

        is_heading = (
            # Markdown headings are unambiguous
            (stripped.startswith("#") and " " in stripped[:7])
            # Title-case standalone lines (not too short, not too long)
            or (
                5 < len(stripped) < 80
                and stripped[0].isupper()
                and stripped == stripped.title()
                and not stripped.endswith((".", ",", ";", ":"))
                and "\t" not in stripped
            )
            # Numbered sections
            or bool(re.match(r"^\d+\.[\d.]*\s+[A-Z]", stripped))
        )

        if is_heading and current_lines:
            # Flush previous section
            content = "\n".join(current_lines).strip()
            if content:
                sections.append(_Section(
                    heading=current_heading,
                    content=content,
                    page_num=page_num,
                ))
            current_heading = stripped.lstrip("#").strip()
            current_lines = []
        elif is_heading:
            current_heading = stripped.lstrip("#").strip()
        else:
            current_lines.append(line)

    # Flush last section
    content = "\n".join(current_lines).strip()
    if content:
        sections.append(_Section(
            heading=current_heading,
            content=content,
            page_num=page_num,
        ))

    # If nothing was detected, return the whole text as one section
    if not sections and text.strip():
        sections.append(_Section(heading="", content=text.strip(), page_num=page_num))

    return sections


# =========================================================================
# Chunking
# =========================================================================

def chunk_text(
    pages: list[tuple[str, int | None]],
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
    filename: str = "",
) -> list[Chunk]:
    """Section-aware chunking with header enrichment.

    1. Detect sections (headings/paragraphs) within each page
    2. Split large sections into overlapping chunks at paragraph/sentence boundaries
    3. Enrich each chunk with filename + section heading for embedding
    """
    if not pages:
        return []

    # Collect all sections across pages
    all_sections: list[_Section] = []
    for page_text, page_num in pages:
        text = page_text.strip()
        if text:
            all_sections.extend(_detect_sections(text, page_num))

    if not all_sections:
        return []

    chunks: list[Chunk] = []
    idx = 0

    for section in all_sections:
        section_chunks = _split_section(
            section.content, chunk_size, chunk_overlap,
        )
        for offset, snippet in section_chunks:
            # Build enriched text: "filename > Section Heading:\ncontent"
            header_parts: list[str] = []
            if filename:
                header_parts.append(filename)
            if section.heading:
                header_parts.append(section.heading)
            header = " > ".join(header_parts)
            enriched = f"{header}:\n{snippet}" if header else snippet

            chunks.append(Chunk(
                index=idx,
                text=snippet,
                enriched_text=enriched,
                page_num=section.page_num,
                char_offset=offset,
                section=section.heading,
            ))
            idx += 1

    return chunks


def chunk_with_parents(
    pages: list[tuple[str, int | None]],
    child_size: int = DEFAULT_CHUNK_SIZE,
    parent_size: int = DEFAULT_PARENT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
    filename: str = "",
) -> tuple[list[Chunk], list[Chunk]]:
    """Create parent-child chunk pairs.

    Returns (child_chunks, parent_chunks) where:
    - child_chunks are small, precise chunks used for vector search
    - parent_chunks are larger context windows returned to the LLM

    Each child_chunk.parent_index points to its parent_chunk.
    """
    # First create parent (large) chunks
    parent_chunks = chunk_text(
        pages, chunk_size=parent_size, chunk_overlap=chunk_overlap, filename=filename,
    )

    # Then create child (small) chunks within each parent
    child_chunks: list[Chunk] = []
    child_idx = 0

    for parent in parent_chunks:
        sub_chunks = _split_section(parent.text, child_size, chunk_overlap)
        for offset, snippet in sub_chunks:
            # Enrich with same header as parent
            header_parts: list[str] = []
            if filename:
                header_parts.append(filename)
            if parent.section:
                header_parts.append(parent.section)
            header = " > ".join(header_parts)
            enriched = f"{header}:\n{snippet}" if header else snippet

            child_chunks.append(Chunk(
                index=child_idx,
                text=snippet,
                enriched_text=enriched,
                page_num=parent.page_num,
                char_offset=parent.char_offset + offset,
                section=parent.section,
                parent_index=parent.index,
            ))
            child_idx += 1

    return child_chunks, parent_chunks


def _split_section(
    text: str,
    chunk_size: int,
    chunk_overlap: int,
) -> list[tuple[int, str]]:
    """Split a section into sized chunks at paragraph/sentence boundaries.

    Returns list of (char_offset, text) tuples.
    """
    text = text.strip()
    if not text:
        return []

    # If it fits in one chunk, return it
    if len(text) <= chunk_size:
        return [(0, text)]

    # Try splitting at paragraph boundaries first
    paragraphs = re.split(r"\n\s*\n", text)
    if len(paragraphs) > 1:
        return _merge_paragraphs(paragraphs, text, chunk_size, chunk_overlap)

    # Fall back to sentence-aware sliding window
    return _sliding_window(text, chunk_size, chunk_overlap)


def _merge_paragraphs(
    paragraphs: list[str],
    original_text: str,
    chunk_size: int,
    chunk_overlap: int,
) -> list[tuple[int, str]]:
    """Merge paragraphs into chunks respecting size limits.

    Keeps whole paragraphs together when possible. Splits oversized
    paragraphs with the sliding window.
    """
    results: list[tuple[int, str]] = []
    current_parts: list[str] = []
    current_len = 0
    running_offset = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # If this single paragraph exceeds chunk_size, split it
        if len(para) > chunk_size:
            # Flush current buffer first
            if current_parts:
                chunk_text_content = "\n\n".join(current_parts)
                offset = original_text.find(current_parts[0], max(running_offset - 10, 0))
                results.append((max(offset, 0), chunk_text_content))
                # Keep last paragraph for overlap
                if chunk_overlap > 0 and current_parts:
                    overlap_text = current_parts[-1]
                    current_parts = [overlap_text] if len(overlap_text) < chunk_overlap else []
                    current_len = sum(len(p) for p in current_parts)
                else:
                    current_parts = []
                    current_len = 0

            # Split the oversized paragraph
            sub_chunks = _sliding_window(para, chunk_size, chunk_overlap)
            para_offset = original_text.find(para, max(running_offset - 10, 0))
            for sub_off, sub_text in sub_chunks:
                results.append((max(para_offset, 0) + sub_off, sub_text))
            running_offset = max(para_offset, 0) + len(para)
            continue

        # Would adding this paragraph exceed the limit?
        added_len = len(para) + (2 if current_parts else 0)  # +2 for "\n\n"
        if current_len + added_len > chunk_size and current_parts:
            # Flush
            chunk_text_content = "\n\n".join(current_parts)
            offset = original_text.find(current_parts[0], max(running_offset - 10, 0))
            results.append((max(offset, 0), chunk_text_content))
            running_offset = max(offset, 0) + len(chunk_text_content)

            # Overlap: keep last paragraph
            if chunk_overlap > 0 and current_parts:
                overlap_text = current_parts[-1]
                current_parts = [overlap_text] if len(overlap_text) < chunk_overlap else []
                current_len = sum(len(p) for p in current_parts)
            else:
                current_parts = []
                current_len = 0

        current_parts.append(para)
        current_len += added_len

    # Flush remaining
    if current_parts:
        chunk_text_content = "\n\n".join(current_parts)
        offset = original_text.find(current_parts[0], max(running_offset - 10, 0))
        results.append((max(offset, 0), chunk_text_content))

    return results


def _sliding_window(
    text: str,
    chunk_size: int,
    chunk_overlap: int,
) -> list[tuple[int, str]]:
    """Sentence-aware sliding window over text.

    Returns list of (offset, snippet) tuples.
    """
    results: list[tuple[int, str]] = []
    offset = 0

    while offset < len(text):
        end = min(offset + chunk_size, len(text))
        snippet = text[offset:end]

        # Break at sentence boundary if not at end of text
        if end < len(text):
            for sep in (". ", "! ", "? ", "\n", "; "):
                last_break = snippet.rfind(sep)
                if last_break > chunk_size // 2:
                    snippet = snippet[: last_break + len(sep)]
                    break

        results.append((offset, snippet.strip()))

        # Done if we reached the end
        if end >= len(text):
            break

        # Advance
        step = max(len(snippet) - chunk_overlap, chunk_size // 2, 1)
        offset += step

    return results
