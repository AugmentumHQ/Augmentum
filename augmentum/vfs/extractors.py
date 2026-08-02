"""Content extraction for the file index.

Wraps the existing `documents.chunker.extract_text` (PDF / DOCX / HTML /
text) and adds the formats it doesn't cover: XLSX, PPTX, EPUB.

The extracted text is what backs the file_index `description` column —
which the enrichment loop then feeds into the embedding model so search
matches the file's contents, not just its filename.

All extractors are best-effort: import failures, parse failures, and
empty results return an empty string rather than raising.  The caller
is the background enrichment loop and shouldn't blow up on a single bad
file.
"""

from __future__ import annotations

import io
import re

from augmentum.utils.logging import get_logger

log = get_logger(__name__)


# --- MIME / extension routing -------------------------------------------

_OOXML_PREFIX = "application/vnd.openxmlformats-officedocument."
_OPENDOC_PREFIX = "application/vnd.oasis.opendocument."

# Extensions we'll route by when MIME is generic (octet-stream / zip).
_EXT_TO_MIME = {
    "pdf":  "application/pdf",
    "docx": _OOXML_PREFIX + "wordprocessingml.document",
    "xlsx": _OOXML_PREFIX + "spreadsheetml.sheet",
    "pptx": _OOXML_PREFIX + "presentationml.presentation",
    "epub": "application/epub+zip",
    "html": "text/html",
    "htm":  "text/html",
    "md":   "text/markdown",
    "txt":  "text/plain",
    "csv":  "text/csv",
    "json": "application/json",
    "xml":  "application/xml",
    "log":  "text/plain",
    "rst":  "text/plain",
    "yaml": "text/plain",
    "yml":  "text/plain",
    "toml": "text/plain",
}


def _resolve_mime(mime: str, filename: str) -> str:
    """Pick the best MIME for routing — fall back to extension when the
    sniffed value is generic (zip envelope, octet-stream, empty)."""
    mime = (mime or "").split(";", 1)[0].strip().lower()
    if mime and mime not in {"application/octet-stream", "application/zip", ""}:
        return mime
    if "." in filename:
        ext = filename.rsplit(".", 1)[-1].lower()
        if ext in _EXT_TO_MIME:
            return _EXT_TO_MIME[ext]
    return mime or "application/octet-stream"


# --- Public API ---------------------------------------------------------

def extract(data: bytes, mime: str, filename: str = "") -> str:
    """Return extracted text for a file's bytes, or "" if no extractor
    applies / all attempts failed.

    The result is a single flat string suitable for storing in the
    file_index.description column.  Callers truncate to their own limit.
    """
    if not data:
        return ""

    resolved = _resolve_mime(mime, filename)

    try:
        if resolved == _OOXML_PREFIX + "spreadsheetml.sheet":
            return _extract_xlsx(data)
        if resolved == _OOXML_PREFIX + "presentationml.presentation":
            return _extract_pptx(data)
        if resolved == "application/epub+zip":
            return _extract_epub(data)

        # Defer to the documents-side extractor for PDF/DOCX/HTML/text.
        # It returns [(text, page_or_none), ...]; we flatten.
        from augmentum.documents.chunker import extract_text as _doc_extract
        pages = _doc_extract(data, resolved, filename=filename)
        return "\n\n".join(p[0] for p in pages if p and p[0])
    except Exception:
        log.debug("vfs_extract_failed", mime=resolved, filename=filename, exc_info=True)
        return ""


def supported_for(mime: str, filename: str = "") -> bool:
    """Cheap pre-flight: is there *any* extractor that would handle this?
    Used by the enrichment loop to skip irrelevant files (audio, video,
    archives, executables) without reading bytes off disk.
    """
    resolved = _resolve_mime(mime, filename)
    if resolved in {
        "application/pdf",
        _OOXML_PREFIX + "wordprocessingml.document",
        _OOXML_PREFIX + "spreadsheetml.sheet",
        _OOXML_PREFIX + "presentationml.presentation",
        "application/msword",
        "application/epub+zip",
        "text/html",
        "application/xhtml+xml",
    }:
        return True
    return resolved.startswith("text/") or resolved in {
        "application/json",
        "application/xml",
    }


# --- XLSX ---------------------------------------------------------------

def _extract_xlsx(data: bytes) -> str:
    """Flatten worksheet cells into a markdown-friendly tabular text blob."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        log.warning("openpyxl_not_installed")
        return ""
    try:
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception:
        log.debug("xlsx_load_failed", exc_info=True)
        return ""

    parts: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        # Cap the per-sheet row scan — we only need enough text to make
        # search hit; spreadsheets with tens of thousands of rows would
        # blow our description budget anyway.
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= 200:
                rows.append("…")
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"# {sheet.title}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


# --- PPTX ---------------------------------------------------------------

def _extract_pptx(data: bytes) -> str:
    """Concatenate slide titles + body text in slide order."""
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python_pptx_not_installed")
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception:
        log.debug("pptx_load_failed", exc_info=True)
        return ""

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text or "" for run in para.runs).strip()
                if line:
                    slide_text.append(line)
        if len(slide_text) > 1:
            parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


# --- EPUB ---------------------------------------------------------------

# EPUB is a ZIP of HTML files. We pull the spine in order, strip tags,
# and concatenate. Avoids depending on `ebooklib` which isn't a current
# requirement.
_TAG_RE = re.compile(r"<[^>]+>")
_WS_RE = re.compile(r"\s+")


def _extract_epub(data: bytes) -> str:
    import zipfile

    from augmentum.utils.safe_archive import UnsafeArchiveError, ensure_zip_sane
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
        ensure_zip_sane(zf, source="vfs_epub")
    except zipfile.BadZipFile:
        log.debug("epub_bad_zip", exc_info=True)
        return ""
    except UnsafeArchiveError as exc:
        log.warning("epub_unsafe_archive", error=str(exc))
        return ""

    try:
        # Read the contained HTML/XHTML files in archive order. The EPUB
        # spec orders the spine via OPF, but for "make this searchable"
        # purposes archive order is good enough and avoids an OPF parser.
        parts: list[str] = []
        for name in zf.namelist():
            lower = name.lower()
            if not (lower.endswith(".xhtml") or lower.endswith(".html") or lower.endswith(".htm")):
                continue
            try:
                raw = zf.read(name).decode("utf-8", errors="replace")
            except Exception as exc:
                log.debug("epub_extractor_read_failed", member=name, error=str(exc))
                continue
            text = _TAG_RE.sub(" ", raw)
            text = _WS_RE.sub(" ", text).strip()
            if text:
                parts.append(text)
            # Cap to avoid runaway extraction on multi-MB books — the
            # description column will truncate further but keep the
            # zip-walk bounded too.
            if sum(len(p) for p in parts) > 200_000:
                break
        return "\n\n".join(parts)
    finally:
        zf.close()
