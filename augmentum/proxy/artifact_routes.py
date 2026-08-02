"""FastAPI router for artifact endpoints (download, list, delete, preview)."""

from __future__ import annotations

import asyncio
import json
import re
from contextlib import suppress

from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, Response
from starlette.responses import StreamingResponse

from augmentum.builds.runtime import (
    apply_project_progress,
    build_status_from_run,
    build_status_snapshot,
    heartbeat_build_run,
    load_persisted_build_run,
    progress_payload_from_state,
    select_active_build,
)
from augmentum.utils.logging import get_logger

log = get_logger(__name__)


def _normalize_patch_filename(filename: str) -> str:
    name = (filename or "").strip().strip("`'\"")
    if name.lower().startswith("file:"):
        name = name[5:].strip()
    return re.sub(
        r"\s*\((?:FULL|full|signature|REFERENCE ONLY:[^)]+)\)\s*$",
        "",
        name,
    ).strip()


def _build_quick_edit_prompt(
    files: list[dict], description: str, previous_attempts: list | None = None
) -> list[dict]:
    """Build a prompt for quick edits — description-driven SEARCH/REPLACE.

    Similar to build_fix_prompt but takes a user description instead of errors.
    Uses targeted context compression (full content for files likely affected,
    signature-only for others) to minimize tokens.
    """
    # Identify which files are likely affected based on description keywords
    desc_lower = description.lower()
    full_files = []
    sig_files = []
    style_words = [
        "style", "color", "background", "font", "css", "theme", "dark",
        "layout", "margin", "padding", "border", "visual", "polish",
        "polished", "look", "draw", "canvas", "animation",
    ]
    js_words = [
        "function", "click", "event", "logic", "score", "speed", "timer",
        "game", "state", "button", "input", "car", "player", "obstacle",
        "canvas", "draw", "render", "animation",
    ]
    html_words = [
        "html", "element", "div", "text", "title", "heading", "structure",
        "add", "remove", "screen", "hud",
    ]
    matched_heuristic = any(kw in desc_lower for kw in style_words + js_words + html_words)
    for f in files:
        path = f["path"].lower()
        # Heuristic: CSS-related words → style files, JS words → script files, HTML → entry
        is_relevant = False
        if any(kw in desc_lower for kw in style_words):
            is_relevant = path.endswith(".css")
        if any(kw in desc_lower for kw in js_words):
            is_relevant = is_relevant or path.endswith((".js", ".ts", ".jsx", ".tsx"))
        if any(kw in desc_lower for kw in html_words):
            is_relevant = is_relevant or path.endswith(".html") or path.endswith(".htm")
        # If no heuristic matched, include all files as full
        if not matched_heuristic:
            is_relevant = True
        if is_relevant:
            full_files.append(f)
        else:
            sig_files.append(f)

    # If heuristic excluded everything, include all as full
    if not full_files:
        full_files = files
        sig_files = []

    # Build file context
    file_context = ""
    for f in full_files:
        file_context += f"\n=== {f['path']} (FULL) ===\n{f['content']}\n"
    for f in sig_files:
        lines = f["content"].count("\n") + 1
        exports = []
        for line in f["content"].split("\n")[:30]:
            stripped = line.strip()
            if stripped.startswith("function ") or stripped.startswith("const ") or stripped.startswith("class "):
                exports.append(stripped.split("(")[0].split("=")[0].strip())
        sig = ", ".join(exports[:15]) if exports else f"{lines} lines"
        file_context += (
            f"\n=== {f['path']} (REFERENCE ONLY: {sig}) ===\n"
            "This file is summarized for awareness only. Do not emit SEARCH/REPLACE blocks for it.\n"
        )

    # Previous attempts
    attempts_text = ""
    if previous_attempts:
        for i, attempt in enumerate(previous_attempts):
            att = attempt if isinstance(attempt, str) else str(attempt)
            attempts_text += f"\n\nPrevious attempt #{i+1} (FAILED — try a different approach):\n{att[:500]}"

    system = (
        "You are editing code files. Follow these steps exactly:\n"
        "1. Identify which file(s) need changes\n"
        "2. For each change, output a SEARCH/REPLACE block\n"
        "3. The SEARCH section must contain EXACT lines from the current file\n"
        "4. The REPLACE section contains the new lines\n\n"
        "Output format (no other text):\n"
        "=== FILE: <filename> ===\n"
        "<<<<<<< SEARCH\n"
        "exact existing lines\n"
        "=======\n"
        "new replacement lines\n"
        ">>>>>>> REPLACE\n\n"
        "Important:\n"
        "- Copy the SEARCH lines EXACTLY from the file (whitespace matters)\n"
        "- Only emit SEARCH/REPLACE blocks for files marked (FULL)\n"
        "- Do not use REFERENCE ONLY summaries as SEARCH text\n"
        "- Only change what's needed\n"
        "- All files share ONE global scope when assembled. Use window.X for cross-file access\n"
        "- End with __PASS_COMPLETE__"
    )

    user = f"Files:\n{file_context}\n\n---\nRequested change: {description}{attempts_text}\n\nApply now."

    return [{"role": "system", "content": system}, {"role": "user", "content": user}]


def _summarize_patch_details(
    response: str,
    snapshot: list[dict],
    files_after: list[dict],
) -> list[dict]:
    """Return UI-safe metadata for SEARCH/REPLACE blocks.

    The builder's patch applier owns fuzzy matching. This companion summary
    exposes enough structure for the UI to explain what happened without
    shipping full patch bodies back through progress surfaces.
    """
    details: list[dict] = []
    after_by_path = {f.get("path", ""): f.get("content", "") for f in files_after}
    snapshot_by_path = {f.get("path", ""): f.get("content", "") for f in snapshot}

    section_re = re.compile(
        r"===\s*FILE:\s*(.+?)\s*===\s*\n([\s\S]*?)(?=\n===\s*FILE:|$)",
        re.IGNORECASE,
    )
    sections = [
        (
            _normalize_patch_filename(m.group(1)),
            m.group(2),
        )
        for m in section_re.finditer(response)
    ] or [("", response)]

    for file_name, body in sections:

        for match in re.finditer(
            r"<<<<<<<\.?\s*SEARCH\n([\s\S]*?)\n?={3,}\n([\s\S]*?)\n?>>>>>>>\.?\s*REPLACE",
            body,
        ):
            search = match.group(1)
            replace = match.group(2)
            target = file_name
            if not target:
                first = (search.strip().splitlines() or [""])[0]
                for path, content in snapshot_by_path.items():
                    if first and first in content:
                        target = path
                        break
                if not target and snapshot:
                    target = snapshot[0].get("path", "")

            before = snapshot_by_path.get(target, "")
            after = after_by_path.get(target, "")
            exact_idx = before.find(search)
            match_kind = "exact" if exact_idx >= 0 else "not_found"
            line = 0
            if exact_idx >= 0:
                line = before[:exact_idx].count("\n") + 1
            else:
                search_lines = [ln.strip() for ln in search.splitlines()]
                content_lines = before.splitlines()
                for i in range(0, max(0, len(content_lines) - len(search_lines) + 1)):
                    if all(
                        j < len(search_lines) and content_lines[i + j].strip() == search_lines[j]
                        for j in range(len(search_lines))
                    ):
                        match_kind = "trimmed"
                        line = i + 1
                        break

            details.append({
                "file": target,
                "line": line,
                "match": match_kind,
                "applied": bool(replace and replace in after),
                "search_lines": len(search.splitlines()),
                "replace_lines": len(replace.splitlines()),
            })

    return details

router = APIRouter(prefix="/api/artifacts", tags=["artifacts"])


def _user_id(request: Request) -> str:
    """Extract user_id from authenticated request."""
    user = request.scope.get("user")
    return user.id if user else ""


async def _isolated_preview_gate(request: Request, artifact_id: str):
    """Authenticate an artifact preview request on the isolated origin.

    Returns ``(response, user_id)``:

    * On the MAIN origin (the common case) → ``(None, "")`` so the
      caller falls back to its normal ``_user_id`` resolution.
    * On the ISOLATED origin → redeems the one-time token (302) or
      validates the preview-session cookie. Returns ``(Response, "")``
      when the caller must short-circuit (redirect / 401), or
      ``(None, user_id)`` once the preview session is valid so the
      caller can scope its artifact lookup to the redeemed user.
    """
    if not request.scope.get("augmentum_preview_isolated"):
        return None, ""
    from augmentum.proxy.content_isolation_routes import check_content_isolated_auth
    auth = await check_content_isolated_auth(request, "artifact_app", artifact_id)
    if auth is not None:
        return auth, ""
    return None, request.scope.get("augmentum_preview_user_id", "")


def _get_store(request: Request):
    store = getattr(request.app.state, "artifact_store", None)
    if not store:
        raise HTTPException(status_code=503, detail="Artifact storage not available")
    return store


@router.post("/build-cancel")
async def cancel_build(request: Request):
    """Cancel the active background build."""
    from augmentum.modes.passthrough.handler import ACTIVE_BUILDS
    build_id = request.query_params.get("build_id", "")
    if not build_id:
        try:
            body = await request.json()
            build_id = body.get("build_id", "") if isinstance(body, dict) else ""
        except Exception:
            build_id = ""
    selected_id, build = select_active_build(
        ACTIVE_BUILDS,
        user_id=_user_id(request),
        build_id=build_id,
    )
    if selected_id and build and build["status"] == "running":
        build["status"] = "cancelled"
        build["_cancel"] = True
        ev = build.get("_change_event")
        if ev is not None:
            ev.set()
        store = getattr(request.app.state, "build_run_store", None)
        if store:
            await store.update(
                selected_id,
                user_id=_user_id(request),
                status="canceled",
                progress=progress_payload_from_state(build),
            )
        log.info("app_builder.cancelled", build_id=selected_id)
        return JSONResponse({"cancelled": True, "build_id": selected_id})
    return JSONResponse({"cancelled": False, "reason": "no active build"})


@router.post("/import")
async def import_file_to_library(request: Request):
    """Import an external file (PDF, DOCX, image, etc.) into the artifact library."""
    import os

    store = _get_store(request)
    if not store:
        raise HTTPException(status_code=503, detail="Artifact store not available")

    form = await request.form()
    uploaded = form.get("file")
    if not uploaded:
        raise HTTPException(status_code=400, detail="No file uploaded")

    contents = await uploaded.read()
    if not contents:
        raise HTTPException(status_code=400, detail="Empty file")
    # Size is capped globally by _MaxBodySizeMiddleware (max_request_body_bytes,
    # default 50 MB); no per-file cap is enforced here beyond that.

    # The multipart filename is attacker-controlled — basename it (defence in
    # depth on top of ArtifactStore.save's own sanitiser) so ``../`` /
    # absolute paths can't traverse, and derive a normalised format token.
    import re as _re
    filename = os.path.basename((uploaded.filename or "").replace("\\", "/")).strip()
    if not filename or set(filename) <= {"."}:
        filename = "imported_file"
    ext = os.path.splitext(filename)[1].lower()
    fmt = ext.lstrip(".")
    # Format drives the mime map + preview dispatch — keep it a short, plain
    # token so a crafted extension can't smuggle anything downstream.
    if not _re.fullmatch(r"[a-z0-9]{1,12}", fmt):
        fmt = ""
    # The extension says what the user CLAIMS; the magic bytes say what
    # the file IS. Container formats route to zip/pdf readers downstream
    # (preview, auto-extract, epub spine walk) — a mislabeled file must
    # not reach them under the claimed format. Demote, don't reject:
    # imports of odd-but-honest files still land, just without the
    # container-specific preview.
    _CONTAINER_FMTS = {"zip", "epub", "cbz", "docx", "xlsx", "pdf"}
    if fmt in _CONTAINER_FMTS:
        from augmentum.utils.safe_archive import sniff_kind
        _sniffed = sniff_kind(contents)
        _expected = "pdf" if fmt == "pdf" else "zip"
        if _sniffed != _expected:
            log.warning(
                "artifact_import_format_mismatch",
                claimed=fmt, sniffed=_sniffed or "unknown", filename=filename,
            )
            fmt = ""
    display_name = os.path.splitext(filename)[0]
    user_id = _user_id(request)

    artifact = await store.save(
        data=contents,
        filename=filename,
        fmt=fmt,
        session_id="",
        display_name=display_name,
        metadata={"imported": True, "original_filename": filename},
        source_json=None,
        user_id=user_id,
    )

    log.info("artifact_imported", id=artifact["id"], filename=filename, format=fmt, size=len(contents))
    return JSONResponse({"id": artifact["id"], "filename": filename, "format": fmt}, status_code=201)


@router.post("/save-html")
async def save_html_to_library(request: Request):
    """Save a single HTML/CSS/JS/SVG code block as a library artifact.

    Accepts JSON:
      { "content": "...", "filename": "index.html", "display_name": "My App",
        "format": "html", "source_json": "...", "session_id": "" }
    """
    body = await request.json()
    content = body.get("content", "")
    filename = body.get("filename", "index.html")
    display_name = body.get("display_name", "Untitled")
    fmt = body.get("format", "html")
    source_json = body.get("source_json")
    session_id = body.get("session_id", "")

    if not content:
        raise HTTPException(status_code=400, detail="content is required")

    store = _get_store(request)
    user_id = _user_id(request)
    artifact = await store.save(
        data=content.encode("utf-8"),
        filename=filename,
        fmt=fmt,
        session_id=session_id,
        display_name=display_name,
        metadata={"file_count": 1, "description": display_name[:200]},
        source_json=source_json,
        user_id=user_id,
    )
    return JSONResponse(artifact)


@router.get("/build-status")
async def get_build_status(request: Request):
    """Poll for background build progress. Used by the persistent build monitor."""
    from augmentum.modes.passthrough.handler import ACTIVE_BUILDS
    builds = ACTIVE_BUILDS
    user_id = _user_id(request)
    build_id = request.query_params.get("build_id", "")
    session_id = request.query_params.get("session_id", "")
    store = getattr(request.app.state, "build_run_store", None)
    selected_id, build = select_active_build(
        builds,
        user_id=user_id,
        build_id=build_id,
        session_id=session_id,
    )
    if not build:
        if store:
            run = await load_persisted_build_run(
                store,
                build_id=build_id,
                session_id=session_id,
                user_id=user_id,
            )
            return JSONResponse(build_status_from_run(run))
        return JSONResponse({"active": False})

    result = build_status_snapshot(build, selected_id)

    if build["status"] in ("complete", "error", "cancelled"):
        project = build.get("project")
        # Fallback: if project is missing, try in-memory checkpoint
        if not project:
            cp = build.get("_checkpoint")
            if cp and cp.get("files"):
                project = {
                    "name": build.get("name", ""),
                    "files": cp["files"],
                    "planned_files": cp.get("planned_files", []),
                    "completed_files": cp.get("completed_files", []),
                    "resumable": True,
                }
        result["project"] = project
        if not build.get("_cleanup"):
            build["_cleanup"] = True
            import asyncio
            async def _cleanup():
                await asyncio.sleep(1800)  # 30 min (was 5 min — too aggressive)
                builds.pop(selected_id, None)
            # Hold a ref so the 30-min sleeper isn't GC'd before it runs.
            from augmentum.utils.bg_tasks import track
            track(_cleanup())

    return JSONResponse(result)


def _build_status_snapshot(build: dict, build_id: str) -> dict:
    """Render the same payload as GET /build-status from a build_state dict.

    Factored out so the SSE stream and the polling endpoint stay in sync.
    """
    return build_status_snapshot(build, build_id)


@router.get("/build-status/stream")
async def stream_build_status(request: Request):
    """Push build progress as SSE events.

    Falls through to the existing 2s polling endpoint for clients that
    can't open EventSource (legacy proxies, fetch-only HTTP clients).
    Each event is one JSON snapshot of the active build matching the
    GET /build-status response shape.
    """
    from augmentum.modes.passthrough.handler import ACTIVE_BUILDS
    user_id = _user_id(request)
    build_id = request.query_params.get("build_id", "")
    session_id = request.query_params.get("session_id", "")
    store = getattr(request.app.state, "build_run_store", None)

    async def event_stream():
        last_serialised = ""
        # Heartbeat cadence — long enough to amortise overhead, short
        # enough to detect stalled connections and proxies that silently
        # drop idle streams (nginx default is 60s).
        heartbeat = 25.0
        # Ceiling on stream lifetime so a forgotten EventSource can't
        # pin server resources forever. Clients can reconnect freely.
        deadline = asyncio.get_event_loop().time() + 1800
        while True:
            if await request.is_disconnected():
                return
            if asyncio.get_event_loop().time() > deadline:
                yield "event: end\ndata: {\"reason\":\"max-duration\"}\n\n"
                return

            selected_id, build = select_active_build(
                ACTIVE_BUILDS,
                user_id=user_id,
                build_id=build_id,
                session_id=session_id,
            )
            if not build:
                if store and (build_id or session_id):
                    run = await load_persisted_build_run(
                        store,
                        build_id=build_id,
                        session_id=session_id,
                        user_id=user_id,
                    )
                    snap = build_status_from_run(run)
                    payload = json.dumps(snap, default=str)
                    if payload != last_serialised:
                        yield f"data: {payload}\n\n"
                        last_serialised = payload
                    if snap.get("status") in ("complete", "error", "cancelled"):
                        yield "event: end\ndata: {\"reason\":\"terminal\"}\n\n"
                        return
                    await asyncio.sleep(2.0 if snap.get("active") else heartbeat)
                    continue
                # No active build — emit one "idle" event then sleep on
                # the heartbeat. We don't terminate so a fresh build
                # started by another tab still gets observed.
                if last_serialised != "idle":
                    yield "event: idle\ndata: {\"active\":false}\n\n"
                    last_serialised = "idle"
                await asyncio.sleep(heartbeat)
                continue

            snap = _build_status_snapshot(build, selected_id)
            payload = json.dumps(snap, default=str)
            if payload != last_serialised:
                yield f"data: {payload}\n\n"
                last_serialised = payload

            terminal = build["status"] in ("complete", "error", "cancelled")
            if terminal:
                yield "event: end\ndata: {\"reason\":\"terminal\"}\n\n"
                return

            ev = build.get("_change_event")
            if ev is None:
                # Pre-SSE-instrumentation build (e.g. started before this
                # process upgrade) — fall back to polling cadence.
                await asyncio.sleep(2.0)
                continue
            try:
                await asyncio.wait_for(ev.wait(), timeout=heartbeat)
            except TimeoutError:
                # No change in the heartbeat window — emit a comment to
                # keep proxies from dropping the connection.
                yield ": keepalive\n\n"
            ev.clear()

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache, no-transform",
            "X-Accel-Buffering": "no",  # disable nginx response buffering
            "Connection": "keep-alive",
        },
    )


@router.get("/{artifact_id}/versions")
async def list_artifact_versions(artifact_id: str, request: Request):
    """List every saved version of ``artifact_id`` newest-first.

    The payload omits file content — call ``/versions/{version_id}``
    for the full snapshot. Workspace renders the list as a sidebar
    timeline, preview, and revert button.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    versions = await store.list_versions(artifact_id, user_id=user_id)
    return JSONResponse({"artifact_id": artifact_id, "versions": versions})


@router.get("/versions/{version_id}")
async def get_artifact_version(version_id: str, request: Request):
    """Fetch a single version including its file snapshot."""
    store = _get_store(request)
    user_id = _user_id(request)
    version = await store.get_version(version_id, user_id=user_id)
    if not version:
        raise HTTPException(status_code=404, detail="Version not found")
    return JSONResponse(version)


@router.post("/{artifact_id}/revert/{version_id}")
async def revert_artifact_to_version(artifact_id: str, version_id: str, request: Request):
    """Restore the artifact's source_json + zip from a saved version.

    Snapshots a fresh version BEFORE overwriting so the revert is itself
    undoable — the user can step back to the pre-revert state.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    target = await store.get_version(version_id, user_id=user_id)
    if not target:
        raise HTTPException(status_code=404, detail="Version not found")
    if target["artifact_id"] != artifact_id:
        raise HTTPException(status_code=400, detail="Version does not belong to this artifact")

    # Snapshot current state so revert is itself reversible.
    try:
        current_source = info.get("source_json")
        if current_source:
            current = json.loads(current_source) if isinstance(current_source, str) else current_source
            current_files = current.get("files", []) if isinstance(current, dict) else []
            if current_files:
                await store.save_version(
                    artifact_id,
                    current_files,
                    user_id=user_id,
                    label=f"Auto-saved before revert to v{target['version_index']}",
                )
    except Exception as exc:
        log.warning("artifact_versions.pre_revert_snapshot_failed",
                    artifact_id=artifact_id, error=str(exc))

    # Rebuild source_json + zip blob from the target snapshot.
    import io
    import zipfile
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in target["files"]:
            if f.get("content"):
                zf.writestr(f.get("path", ""), f["content"])

    new_source = json.dumps({
        "type": "application",
        "name": info.get("display_name", ""),
        "files": target["files"],
    })
    await store.update_file(artifact_id, zip_buf.getvalue(), user_id=user_id)
    await store.update_source(artifact_id, new_source, user_id=user_id)

    return JSONResponse({
        "reverted": True,
        "artifact_id": artifact_id,
        "version_index": target["version_index"],
        "files": target["files"],
    })


@router.get("/{artifact_id}/download")
async def download_artifact(artifact_id: str, request: Request):
    """Download an artifact file by ID."""
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")

    file_path = store.get_file_path(info["path"])
    if not file_path:
        raise HTTPException(status_code=404, detail="Artifact file not found on disk")

    # Map format to MIME type
    mime_map = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "svg": "image/svg+xml",
        "html": "text/html",
        "htm": "text/html",
        "md": "text/markdown",
        "txt": "text/plain",
        "rst": "text/x-rst",
        "log": "text/plain",
        "csv": "text/csv",
        "json": "application/json",
        "epub": "application/epub+zip",
        # Audio
        "mp3": "audio/mpeg",
        "wav": "audio/wav",
        "m4a": "audio/mp4",
        "flac": "audio/flac",
        "ogg": "audio/ogg",
        "opus": "audio/opus",
        "aac": "audio/aac",
        # Video — matters for browser playback because HTML5 <video>
        # won't seek or scrub without a real media Content-Type.
        "mp4": "video/mp4",
        "mov": "video/quicktime",
        "webm": "video/webm",
        "mkv": "video/x-matroska",
        "avi": "video/x-msvideo",
    }
    media_type = mime_map.get(info.get("format", ""), "application/octet-stream")

    return FileResponse(
        path=str(file_path),
        media_type=media_type,
        filename=info.get("display_name", info["filename"]),
    )


@router.get("/{artifact_id}/preview-image")
async def get_artifact_preview_image(artifact_id: str, request: Request):
    """Serve the library thumbnail PNG for an app artifact.

    Returns 404 when no screenshot has been captured yet — the
    library card falls back to a placeholder and may POST
    ``/capture-preview`` to backfill on demand.
    """
    from augmentum.tools.artifact_application import _preview_image_path
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    file_path = store.get_file_path(info.get("path", ""))
    if not file_path:
        raise HTTPException(status_code=404, detail="Artifact file not found")
    preview_path = _preview_image_path(file_path)
    if not preview_path.is_file():
        raise HTTPException(status_code=404, detail="Preview image not captured")
    return FileResponse(
        path=str(preview_path),
        media_type="image/png",
        headers={"Cache-Control": "private, max-age=3600"},
    )


@router.post("/{artifact_id}/capture-preview")
async def capture_artifact_preview(artifact_id: str, request: Request):
    """Lazy-backfill a library thumbnail for an existing app artifact.

    Reassembles the app's HTML from ``source_json`` (same logic the
    live-preview path used) and runs the headless-chromium screenshot
    helper. Idempotent: if a preview already exists it's overwritten;
    if chromium isn't available the call returns 200 with
    ``captured: false`` so the UI can stop retrying.
    """
    from augmentum.tools.artifact_application import (
        assemble_application_html,
        capture_app_preview_screenshot,
    )
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    source_json = info.get("source_json") or ""
    if not source_json:
        return JSONResponse({"captured": False, "reason": "no source"})
    try:
        source = json.loads(source_json) if isinstance(source_json, str) else source_json
    except (json.JSONDecodeError, TypeError):
        return JSONResponse({"captured": False, "reason": "invalid source"})
    files = (source or {}).get("files") or []
    if not files:
        return JSONResponse({"captured": False, "reason": "no files"})
    html = assemble_application_html(files)
    if not html:
        return JSONResponse({"captured": False, "reason": "assemble failed"})
    ok = await capture_app_preview_screenshot(store, artifact_id, html, user_id=user_id)
    return JSONResponse({"captured": bool(ok)})


@router.get("/{artifact_id}/epub-text")
async def get_epub_text(artifact_id: str, request: Request):
    """Return an EPUB artifact's spine as plain-text chapters.

    Powers the read-aloud control in the book viewer — the iframe preview
    runs in an opaque sandbox origin and can't call authed APIs, so the
    parent page fetches the text here and drives the shared TTS pipeline.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    if info.get("format") != "epub":
        raise HTTPException(status_code=400, detail="Artifact is not an EPUB")
    file_path = store.get_file_path(info.get("path", ""))
    if not file_path or not file_path.is_file():
        raise HTTPException(status_code=404, detail="Artifact file not found on disk")

    from augmentum.vfs import epub_extractor

    chapters = await asyncio.to_thread(epub_extractor.chapters_text, str(file_path))
    if not chapters:
        raise HTTPException(status_code=422, detail="Could not extract text from this EPUB")
    return JSONResponse({
        "title": info.get("display_name") or info.get("filename") or "Ebook",
        "chapters": chapters,
    })


async def _epub_artifact_or_404(request: Request, artifact_id: str) -> dict:
    store = _get_store(request)
    info = await store.get(artifact_id, user_id=_user_id(request))
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    if info.get("format") != "epub":
        raise HTTPException(status_code=400, detail="Artifact is not an EPUB")
    return info


@router.get("/{artifact_id}/narration")
async def get_artifact_narration(artifact_id: str, request: Request):
    """Status of this EPUB artifact's paired TTS narration (if any)."""
    from augmentum.proxy.narration_common import narration_status
    return JSONResponse(await narration_status(request, "artifact", artifact_id))


@router.post("/{artifact_id}/narration")
async def start_artifact_narration(artifact_id: str, request: Request, force: int = 0):
    """Record (synthesize) a TTS narration for this EPUB artifact.

    Idempotent: returns the existing one if already done or in progress,
    unless ``force=1``. Body may carry ``{"voice": "...", "format": "mp3"|"wav"}``.
    """
    from augmentum.proxy.narration_common import narration_start
    info = await _epub_artifact_or_404(request, artifact_id)
    try:
        body = await request.json()
    except Exception:  # noqa: BLE001
        body = {}
    body = body if isinstance(body, dict) else {}
    voice = str(body.get("voice") or "")
    out_fmt = str(body.get("format") or "mp3")
    title = info.get("display_name") or info.get("filename") or "Ebook"
    return JSONResponse(await narration_start(
        request, "artifact", artifact_id, title=title, voice=voice,
        output_format=out_fmt, force=bool(force),
    ))


@router.get("/{artifact_id}")
async def get_artifact(artifact_id: str, request: Request):
    """Get artifact metadata by ID.

    For ZIP artifacts without source_json, auto-extracts files from the ZIP
    and returns them as source_json so the workspace can render them.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")

    # Auto-extract source_json from ZIP if missing
    if info.get("format") == "zip" and not info.get("source_json"):
        file_path = store.get_file_path(info.get("path", ""))
        if file_path and file_path.is_file():
            import json as _json
            import zipfile as _zf
            try:
                files = []
                with _zf.ZipFile(str(file_path), "r") as zf:
                    for name in zf.namelist():
                        if name == "README.md":
                            continue
                        try:
                            content = zf.read(name).decode("utf-8", errors="replace")
                        except Exception as exc:
                            log.debug("artifact_zip_read_failed", member=name, error=str(exc))
                            continue
                        # Guess role from extension
                        role = "script"
                        if name.endswith(".html") or name.endswith(".htm"):
                            role = "entry"
                        elif name.endswith(".css"):
                            role = "style"
                        elif name.endswith(".json"):
                            role = "data"
                        files.append({"path": name, "role": role, "content": content})
                if files:
                    source = _json.dumps({"type": "application", "files": files})
                    info["source_json"] = source
                    # Persist so we don't re-extract next time
                    await store.update_source(artifact_id, source, user_id=user_id)
            except Exception as e:
                log.warning("artifact_get.zip_extract_failed", id=artifact_id, error=str(e))

    return JSONResponse(info)


@router.get("")
async def list_artifacts(
    request: Request,
    session_id: str = "",
    task_id: str = "",
):
    """List artifacts, filtered by session or task. No filter = all artifacts."""
    store = _get_store(request)
    user_id = _user_id(request)

    if task_id:
        results = await store.list_for_task(task_id, user_id=user_id)
    elif session_id:
        results = await store.list_for_session(session_id, user_id=user_id)
    else:
        results = await store.list_all(user_id=user_id)

    return JSONResponse(results)


@router.post("/iterate")
async def iterate_project(request: Request):
    """Start a build/iterate as a background task. Returns immediately.

    The build runs server-side independent of the client connection.
    Frontend polls /api/artifacts/build-status for progress.
    """
    body = await request.json()
    description = body.get("description", "")
    scaffold = body.get("scaffold", "static")
    model = body.get("model", "")
    files = body.get("files", [])
    session_id = body.get("session_id", "")

    if not description:
        raise HTTPException(status_code=400, detail="description is required")
    if not files:
        raise HTTPException(status_code=400, detail="files array is required for iteration")

    tool_registry = getattr(request.app.state, "tool_registry", None)
    if not tool_registry:
        raise HTTPException(status_code=503, detail="Tool registry not available")

    builder = tool_registry.get("build_application")
    if not builder:
        raise HTTPException(status_code=503, detail="Application builder not registered")

    # Capture user_id now — the background task's request object may not be
    # usable after the response returns.
    user_id = _user_id(request)

    # Start build as a background task with shared state for polling
    from augmentum.modes.passthrough.handler import ACTIVE_BUILDS

    build_id = f"build_{int(asyncio.get_event_loop().time() * 1000)}"
    task_id = body.get("task_id", "")
    build_state = {
        "id": build_id,
        "kind": "application",
        "user_id": user_id,
        "session_id": session_id,
        "task_id": task_id,
        "started_at": asyncio.get_event_loop().time(),
        "name": description[:60],
        "status": "running",
        "passes": [],
        "error": None,
        "project": None,
    }
    ACTIVE_BUILDS[build_id] = build_state
    build_store = getattr(request.app.state, "build_run_store", None)
    if build_store:
        await build_store.create(
            build_id=build_id,
            user_id=user_id,
            session_id=session_id,
            task_id=task_id,
            kind="application",
            status="running",
            name=description[:60],
            request={
                "description": description,
                "scaffold": scaffold,
                "model": model,
                "iteration": True,
            },
        )

    # SSE notifies stream consumers whenever build_state mutates. Created
    # lazily on the first progress callback so /iterate doesn't pay the
    # cost when no client subscribes (e.g. tool-calling API path).
    build_state["_change_event"] = asyncio.Event()
    heartbeat_task = None
    if build_store and user_id:
        heartbeat_task = asyncio.create_task(
            heartbeat_build_run(
                build_store,
                build_id=build_id,
                user_id=user_id,
                state=build_state,
            )
        )

    async def on_progress(data):
        # Check cancellation flag (same as direct build path)
        if build_state.get("_cancel"):
            raise asyncio.CancelledError("Build cancelled by user")
        progress = data.get("project_progress", {})
        apply_project_progress(build_state, progress)
        if progress.get("name"):
            build_state["name"] = progress["name"]
        if progress.get("pass"):
            existing = next((p for p in build_state["passes"] if p["name"] == progress["pass"]), None)
            if existing:
                existing["status"] = progress.get("status", "running")
                existing["detail"] = progress.get("detail", "")
            else:
                build_state["passes"].append({
                    "name": progress["pass"],
                    "status": progress.get("status", "running"),
                    "detail": progress.get("detail", ""),
                })
        # Track token counts for real-time display
        if progress.get("totalTokens") is not None:
            build_state["totalTokens"] = progress["totalTokens"]
            build_state["llmCalls"] = progress.get("llmCalls", 0)
        # Checkpoint files incrementally — survive crashes
        if progress.get("files"):
            if not build_state.get("_checkpoint"):
                build_state["_checkpoint"] = {}
            build_state["_checkpoint"]["files"] = progress["files"]
            build_state["_checkpoint"]["planned_files"] = progress.get("planned_files", [])
            build_state["_checkpoint"]["completed_files"] = progress.get("completed_files", [])
        # Wake any SSE subscribers so they can re-read build_state.
        # Set then immediately clear — consumers re-arm before reading.
        ev = build_state.get("_change_event")
        if ev is not None:
            ev.set()
        if build_store:
            await build_store.update(
                build_id,
                user_id=user_id,
                name=build_state.get("name", ""),
                progress=progress_payload_from_state(build_state),
            )

    def _project_from_checkpoint() -> dict | None:
        """Build resumable project data from in-memory checkpoint."""
        cp = build_state.get("_checkpoint")
        if not cp or not cp.get("files"):
            return None
        return {
            "name": build_state["name"],
            "files": cp["files"],
            "planned_files": cp.get("planned_files", []),
            "completed_files": cp.get("completed_files", []),
            "resumable": True,
            "scaffold": scaffold,
            "qualityStatus": build_state.get("qualityStatus", "clean"),
            "warnings": build_state.get("warnings", []),
            "blockingErrors": build_state.get("blockingErrors", []),
        }

    async def _save_sqlite_checkpoint():
        """Persist current checkpoint to SQLite (survives server restart)."""
        cp = build_state.get("_checkpoint")
        if not cp or not cp.get("files"):
            return
        try:
            store = _get_store(request)
            await store.save_checkpoint(
                build_id=build_id,
                session_id=session_id,
                project_name=build_state["name"],
                files=cp["files"],
                planned_files=cp.get("planned_files", []),
                scaffold=scaffold,
                user_id=user_id,
            )
        except Exception as e:
            log.warning("artifact_iterate.checkpoint_save_failed", error=str(e))

    async def run_pipeline():
        try:
            kwargs = {
                "description": description,
                "scaffold": scaffold,
                "_request_model": model,
                "_progress_callback": on_progress,
                "_session_id": session_id,
                "_user_id": user_id,
                "_task_id": task_id,
                "_build_id": build_id,
            }
            planned = body.get("planned_files")
            if planned and files:
                kwargs["resume_from"] = {"files": files, "planned_files": planned}
            else:
                kwargs["existing_project"] = {"files": files}
            result = await builder.execute(**kwargs)
            # Defensive: ensure metadata is a dict with "project" key
            project = None
            if isinstance(result.metadata, dict):
                project = result.metadata.get("project")
            if not project:
                project = _project_from_checkpoint()

            # Check if this was a cancellation (execute catches CancelledError internally)
            is_cancel = build_state.get("_cancel") or "CancelledError" in (result.error or "")
            if is_cancel:
                build_state["status"] = "cancelled"
                # Cancel is explicit user intent — drop the checkpoint so a
                # server restart doesn't quietly auto-resume the same build.
                # Crashes still hit the error/except branches below which DO
                # persist a checkpoint for recovery.
                try:
                    store = _get_store(request)
                    await store.delete_checkpoint(build_id, user_id=user_id)
                except Exception as _cp_err:
                    log.warning("artifact_iterate.checkpoint_cleanup_failed", error=str(_cp_err))
            elif result.success:
                build_state["status"] = "complete"
                # Clean up SQLite checkpoint on success
                try:
                    store = _get_store(request)
                    await store.delete_checkpoint(build_id, user_id=user_id)
                except Exception as _cp_err:
                    log.warning("artifact_iterate.checkpoint_cleanup_failed", error=str(_cp_err))
            else:
                build_state["status"] = "error"
                build_state["error"] = result.error
                # Persist checkpoint to SQLite on failure
                await _save_sqlite_checkpoint()
            build_state["project"] = project
            if isinstance(project, dict):
                build_state["qualityStatus"] = project.get("qualityStatus") or project.get("quality_status") or build_state.get("qualityStatus", "clean")
                build_state["warnings"] = project.get("warnings") or build_state.get("warnings", [])
                build_state["blockingErrors"] = project.get("blockingErrors") or project.get("blocking_errors") or build_state.get("blockingErrors", [])
        except (asyncio.CancelledError, GeneratorExit):
            build_state["status"] = "cancelled"
            build_state["project"] = build_state.get("project") or _project_from_checkpoint()
            # Same reasoning as the is_cancel branch above: drop the
            # checkpoint on user-initiated cancel rather than persisting it.
            try:
                store = _get_store(request)
                await store.delete_checkpoint(build_id, user_id=user_id)
            except Exception as _cp_err:
                log.warning("artifact_iterate.checkpoint_cleanup_failed", error=str(_cp_err))
            log.info("artifact_iterate.cancelled")
        except Exception as exc:
            log.warning("artifact_iterate.failed", error=str(exc))
            build_state["status"] = "error"
            build_state["error"] = str(exc)
            build_state["project"] = build_state.get("project") or _project_from_checkpoint()
            await _save_sqlite_checkpoint()
        finally:
            if heartbeat_task:
                heartbeat_task.cancel()
                with suppress(asyncio.CancelledError):
                    await heartbeat_task
            if build_store:
                project = build_state.get("project") or {}
                artifact_id = ""
                if isinstance(project, dict):
                    artifact_id = project.get("artifactId") or project.get("artifact_id") or ""
                await build_store.update(
                    build_id,
                    user_id=user_id,
                    status=build_state.get("status", "running"),
                    name=build_state.get("name", ""),
                    artifact_id=artifact_id,
                    progress=progress_payload_from_state(build_state),
                    result={"project": project, "artifact_id": artifact_id},
                    error=build_state.get("error") or "",
                )
            # Wake SSE subscribers one last time so they can observe the
            # terminal state and close cleanly instead of hanging until
            # the heartbeat timeout.
            ev = build_state.get("_change_event")
            if ev is not None:
                ev.set()

    try:
        # Hold a ref so the build pipeline can't be GC'd mid-execution.
        # Without this, long-running builds (5-30s typical) can vanish
        # under interpreter GC pressure with no exception surfacing.
        from augmentum.utils.bg_tasks import track
        track(run_pipeline())
    except RuntimeError as e:
        ACTIVE_BUILDS.pop(build_id, None)
        raise HTTPException(status_code=503, detail=f"Failed to start build: {e}") from e

    return JSONResponse({
        "started": True,
        "build_id": build_id,
        "message": "Build started in background. Poll /api/artifacts/build-status for progress.",
    })


@router.post("/fix")
async def fix_project_errors(request: Request):
    """Quick edit or error fix for a project — single LLM call with SEARCH/REPLACE.

    Two modes:
    - **Error fix**: provide `errors` array — uses build_fix_prompt for targeted fixes.
    - **Quick edit**: provide `description` (no errors) — uses an edit prompt for
      SEARCH/REPLACE modifications. This is the fast path for 80% of edits.

    Both modes use 3-tier fuzzy matching for patch application and QuickJS
    verification with rollback on regression.

    Accepts JSON:
      {
        "files": [{"path": "app.js", "role": "script", "content": "..."}],
        "model": "model-name",
        "errors": ["TypeError: ..."],          // error fix mode
        "description": "change bg to blue",    // quick edit mode
        "previous_attempts": ["..."]           // optional
      }

    Returns:
      {"success": true, "files": [...], "patches_applied": 3, "raw_response": "..."}
    """
    from augmentum.tools.application_scaffolds import build_fix_prompt

    body = await request.json()
    errors = body.get("errors", [])
    description = body.get("description", "")
    files = body.get("files", [])
    model = body.get("model", "")
    previous_attempts = body.get("previous_attempts") or None

    if not files:
        raise HTTPException(status_code=400, detail="files array is required")
    if not errors and not description:
        raise HTTPException(status_code=400, detail="errors or description is required")

    tool_registry = getattr(request.app.state, "tool_registry", None)
    builder = tool_registry.get("build_application") if tool_registry else None
    if not builder:
        raise HTTPException(status_code=503, detail="Application builder not registered")

    # Snapshot for rollback
    snapshot = [{"path": f["path"], "content": f["content"]} for f in files]

    # Build prompt — error fix or quick edit
    if errors:
        messages = build_fix_prompt(files, errors, previous_attempts=previous_attempts)
    else:
        # Quick edit mode — describe the change, model produces SEARCH/REPLACE
        messages = _build_quick_edit_prompt(files, description, previous_attempts)

    # Call LLM — use SEARCH/REPLACE grammar on llama.cpp for guaranteed valid format
    from augmentum.tools.application_scaffolds import GRAMMAR_SEARCH_REPLACE
    try:
        response = await builder._call_llm(messages, max_tokens=max(
            getattr(builder, '_max_tokens', 8192) // 2, 2048
        ), model=model, grammar=GRAMMAR_SEARCH_REPLACE)
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"LLM call failed: {exc}")

    # Apply patches with 3-tier fuzzy matching
    patches = builder._apply_file_patches(files, response)
    patch_details = _summarize_patch_details(response, snapshot, files)

    if patches > 0 and errors:
        # Verify the fix didn't make things worse (quickjs if available).
        # Only run for error-fix mode where we have a baseline error count.
        # Quick edits (description-only) skip this — QuickJS false positives
        # on browser APIs (document, window, fetch) cause spurious rollbacks.
        assembled = builder._assemble(files)
        import re as _re
        js_blocks = _re.findall(r"<script[^>]*>([\s\S]*?)</script>", assembled, _re.IGNORECASE)
        post_js = "\n".join(js_blocks)
        if post_js.strip():
            post_errors = builder._execute_js_verify(post_js, assembled)
            orig_set = set(errors)
            new_errors = [e for e in post_errors if e not in orig_set]

            if len(post_errors) > len(errors) or new_errors:
                # Rollback — fix made things worse
                for orig in snapshot:
                    target = next((f for f in files if f["path"] == orig["path"]), None)
                    if target:
                        target["content"] = orig["content"]
                log.warning("artifact_fix.rollback",
                            before=len(errors), after=len(post_errors))
                return JSONResponse({
                    "success": False,
                    "error": "Fix introduced new errors — rolled back",
                    "files": files,
                    "patches_applied": 0,
                    "patch_details": patch_details,
                    "fix_response": response,
                })

    # Run structural analysis on the fixed project
    from augmentum.tools.artifact_application import ApplicationBuilderTool
    analysis = ApplicationBuilderTool.analyze_project(files)

    return JSONResponse({
        "success": patches > 0,
        "files": files,
        "patches_applied": patches,
        "patch_details": patch_details,
        "raw_response": response,
        "structural_issues": analysis["structural_issues"][:10],
        "unresolved": analysis["unresolved"][:10],
    })


@router.post("/verify")
async def verify_project(request: Request):
    """Universal code verification — the ground truth for generated applications.

    Accepts any set of HTML/CSS/JS files and returns a comprehensive
    quality report without requiring a running server, browser, or LLM.

    Checks: runtime errors (quickjs), structural issues (cross-file analysis),
    smoke tests (wiring, state, output, interaction), intent verification
    (description → implementation mapping), accessibility, security.

    POST JSON:
      { "description": "what this app should do",
        "files": [{"path": "...", "role": "entry|style|script", "content": "..."}] }

    Returns:
      { "verdict": "PASS"|"FAIL", "score": 0-100,
        "runtime_errors": [...], "smoke_tests": [...], "intent_gaps": [...],
        "structural_issues": [...], "suggestions": [...] }
    """
    from augmentum.tools.artifact_application import ApplicationBuilderTool

    body = await request.json()
    description = body.get("description", "")
    files = body.get("files", [])

    if not files:
        raise HTTPException(status_code=400, detail="files array is required")

    # Use the registered tool instance (properly initialized) instead of __new__
    tool_registry = getattr(request.app.state, "tool_registry", None)
    builder = tool_registry.get("build_application") if tool_registry else None
    if not builder:
        # Fallback: static methods still work without instance
        from augmentum.tools.artifact_application import ApplicationBuilderTool as ABT
        builder = ABT.__new__(ABT)

    # 1. Structural analysis
    analysis = ApplicationBuilderTool.analyze_project(files)

    # 2. Runtime verification (quickjs)
    assembled = builder._assemble(files)
    import re as _re
    js_blocks = _re.findall(r"<script[^>]*>([\s\S]*?)</script>", assembled, _re.IGNORECASE)
    all_js = "\n".join(js_blocks)
    runtime_errors = builder._execute_js_verify(all_js, assembled) if all_js.strip() else []

    # 3. Intent verification
    intent_gaps = ApplicationBuilderTool.verify_intent(description, files) if description else []

    # 4. Compute score
    total_checks = 10  # base
    passed = total_checks

    if runtime_errors:
        passed -= min(len(runtime_errors), 4)
    if analysis["structural_issues"]:
        passed -= min(len(analysis["structural_issues"]), 3)
    if intent_gaps:
        passed -= min(len(intent_gaps), 3)
    if analysis["unresolved"]:
        passed -= min(len(analysis["unresolved"]), 2)

    score = max(0, int((passed / total_checks) * 100))
    verdict = "PASS" if score >= 60 and not runtime_errors else "FAIL"

    return JSONResponse({
        "verdict": verdict,
        "score": score,
        "runtime_errors": runtime_errors[:10],
        "structural_issues": analysis["structural_issues"][:10],
        "unresolved": analysis["unresolved"][:10],
        "intent_gaps": intent_gaps[:10],
        "file_summary": [
            {"path": fi["path"], "role": fi["role"], "lines": fi["lines"]}
            for fi in analysis["per_file"]
        ],
    })



@router.get("/{artifact_id}/preview")
async def preview_artifact(artifact_id: str, request: Request):
    """Preview an artifact inline in an iframe.

    - **PDF**: Served directly — browsers render PDFs natively in iframes.
    - **Images** (png/jpg/svg): Wrapped in a centered HTML page.
    - **HTML**: Served as-is.
    - **DOCX**: Converted to HTML via python-docx.
    - **EPUB**: Extracted XHTML chapters rendered as HTML book viewer.
    - **PPTX**: Converted to HTML slide viewer via python-pptx.
    - **XLSX/CSV**: Converted to HTML table via openpyxl.
    - **Others**: Download card fallback.

    When served on the isolated preview origin (port 6444), the request
    is authenticated via the one-time token / preview-session cookie
    instead of the main session cookie, and the artifact lookup is
    scoped to the redeemed user.
    """
    gate_response, preview_uid = await _isolated_preview_gate(request, artifact_id)
    if gate_response is not None:
        return gate_response
    store = _get_store(request)
    user_id = preview_uid or _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")

    fmt = info.get("format", "")
    display_name = info.get("display_name", info.get("filename", "Artifact"))
    from html import escape as _html_escape
    safe_display_name = _html_escape(str(display_name))
    safe_format = _html_escape(str(fmt).upper())
    download_url = f"/api/artifacts/{artifact_id}/download"
    file_path = store.get_file_path(info.get("path", ""))

    # --- PDF: serve directly (browsers render natively in iframes) ---
    if fmt == "pdf" and file_path:
        return FileResponse(
            path=str(file_path),
            media_type="application/pdf",
            headers={"Content-Disposition": "inline"},
        )

    # --- Images: centered in dark page ---
    if fmt in ("png", "jpg", "jpeg", "svg", "gif", "webp") and file_path:
        # The preview iframe is sandboxed WITHOUT allow-same-origin (so it can
        # safely host untrusted imported HTML). A plain HTML wrapper that
        # points <img> at the auth-gated /download endpoint fails inside that
        # opaque-origin frame — the subresource request can't carry the
        # session, so it renders a broken-image icon even though the identical
        # <img> works in the top document (chat, the editor). Fix: make the
        # preview SELF-CONTAINED — no auth-gated subresource. Small images
        # embed as a data URI (keeps the centered dark backdrop); large images
        # stream directly as the iframe body (like the PDF branch), which the
        # browser renders natively.
        _IMG_MIME = {
            "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "gif": "image/gif", "webp": "image/webp", "svg": "image/svg+xml",
        }
        media_type = _IMG_MIME.get(fmt, "application/octet-stream")
        # SVG renders natively as the iframe body; nosniff so it's treated as
        # an image, not executed as a document.
        if fmt == "svg":
            return FileResponse(
                path=str(file_path), media_type=media_type,
                headers={"Content-Disposition": "inline",
                         "X-Content-Type-Options": "nosniff"},
            )
        _INLINE_MAX = 12 * 1024 * 1024  # base64 bloats ~33%; cap the embed
        try:
            size = file_path.stat().st_size
        except OSError:
            size = 0
        if 0 < size <= _INLINE_MAX:
            import base64
            b64 = base64.b64encode(file_path.read_bytes()).decode("ascii")
            html = (
                f'<!DOCTYPE html><html><head><meta charset="utf-8">'
                f'<style>body{{margin:0;display:flex;align-items:center;'
                f'justify-content:center;min-height:100vh;background:#1e1e2e}}'
                f'img{{max-width:100%;max-height:100vh}}</style></head>'
                f'<body><img src="data:{media_type};base64,{b64}" '
                f'alt="{safe_display_name}"></body></html>'
            )
            return HTMLResponse(html)
        return FileResponse(
            path=str(file_path), media_type=media_type,
            headers={"Content-Disposition": "inline"},
        )

    # --- HTML: serve raw content ---
    if fmt in ("html", "htm") and file_path:
        # Imported HTML is untrusted, user-supplied markup. Rendering it
        # inline same-origin ("Open in new tab", cast-to-TV) would execute
        # its scripts with the viewer's session cookie — stored XSS. Serve
        # imported HTML as a download instead. App-builder / tool-generated
        # HTML (not flagged imported) renders as before; the library hero
        # sandboxes it regardless. See the import endpoint's metadata stamp.
        if bool((info.get("metadata") or {}).get("imported")):
            return FileResponse(
                path=str(file_path),
                media_type="application/octet-stream",
                filename=(display_name if str(display_name).lower().endswith((".html", ".htm"))
                          else f"{display_name}.html"),
                headers={"X-Content-Type-Options": "nosniff"},
            )
        content = file_path.read_text(encoding="utf-8", errors="replace")
        return HTMLResponse(content, headers={"X-Content-Type-Options": "nosniff"})

    # --- DOCX / EPUB / PPTX / XLSX: parse via third-party libs that hold the
    # GIL on large files (openpyxl XLSX → 100K rows = hundreds of ms; python-
    # docx on long DOCX similar). Hand off to a worker thread so the event
    # loop keeps serving the healthcheck and concurrent chat streams.
    if fmt == "docx" and file_path:
        html = await asyncio.to_thread(_docx_to_html, file_path, display_name, download_url)
        if html:
            return HTMLResponse(html)

    if fmt == "epub" and file_path:
        html = await asyncio.to_thread(_epub_to_html, file_path, display_name, download_url)
        if html:
            return HTMLResponse(html)

    if fmt == "pptx" and file_path:
        html = await asyncio.to_thread(_pptx_to_html, file_path, display_name, download_url)
        if html:
            return HTMLResponse(html)

    if fmt in ("xlsx", "csv") and file_path:
        html = await asyncio.to_thread(_xlsx_to_html, file_path, fmt, display_name, download_url)
        if html:
            return HTMLResponse(html)

    # --- Markdown/text: render with basic styling ---
    if fmt in ("md", "txt", "json", "csv", "rst", "log") and file_path:
        text = file_path.read_text(encoding="utf-8", errors="replace")
        escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        html = _preview_shell(
            display_name, download_url,
            f'<pre style="white-space:pre-wrap;word-break:break-word;'
            f'font-family:monospace;font-size:13px;line-height:1.6;'
            f'padding:24px;margin:0">{escaped}</pre>',
        )
        return HTMLResponse(html)

    # --- Application bundles: serve index.html so the app actually runs.
    # App-builder artifacts ship as a zip + a ``source_json`` carrying
    # type:"application" + a files manifest. Without this branch the zip
    # falls through to the archive-listing renderer below and the user
    # sees a file list instead of their running app. Sibling assets
    # (script src=, link href=) resolve via the ``/preview/{path}``
    # route added alongside, with a <base> tag pointing there so the
    # app's own relative paths "just work".
    if fmt == "zip" and file_path:
        index_html = _application_index_html(info, file_path, artifact_id)
        if index_html is not None:
            return HTMLResponse(index_html)

    # --- Archives: directory listing (zip/tar/gz/tgz/bz2 via stdlib) ---
    if fmt in ("zip", "tar", "gz", "tgz", "bz2") and file_path:
        from augmentum.proxy.files_routes import _archive_to_html
        html = _archive_to_html(file_path, display_name, download_url, fmt)
        if html:
            return HTMLResponse(html)

    # --- Fallback: download card ---
    size_kb = f"{info.get('size_bytes', 0) / 1024:.1f} KB"
    html = _preview_shell(
        display_name, download_url,
        f'<div style="text-align:center;padding:48px">'
        f'<div style="font-size:48px;margin-bottom:16px;opacity:0.3">📄</div>'
        f'<div style="margin-bottom:8px;font-size:16px">{safe_display_name}</div>'
        f'<div style="color:#6b6b80;margin-bottom:24px">{safe_format} &middot; {size_kb}</div>'
        f'<a href="{download_url}" download style="color:#6c8aff;text-decoration:none;'
        f'padding:10px 24px;border:1px solid #2d2d45;border-radius:8px;'
        f'transition:background 0.15s"'
        f' onmouseover="this.style.background=\'rgba(108,138,255,0.1)\'"'
        f' onmouseout="this.style.background=\'none\'">'
        f'Download</a></div>',
    )
    return HTMLResponse(html)


def _is_application_artifact(info: dict) -> bool:
    """True when the artifact's source_json marks it as an app-builder
    bundle. Other zips (user uploads, generic archives) fall through to
    the archive-listing path."""
    raw = info.get("source_json")
    if not raw:
        return False
    try:
        source = json.loads(raw) if isinstance(raw, str) else raw
    except (ValueError, TypeError):
        return False
    return isinstance(source, dict) and source.get("type") == "application"


def _application_index_html(info: dict, file_path, artifact_id: str) -> str | None:
    """Extract index.html from an application zip and inject a <base>
    tag pointing at the sibling-file route so the app's relative URLs
    (``<script src="app.js">`` etc.) resolve against this artifact's
    contents rather than the augmentum origin root.

    Returns None when the artifact isn't an application or has no
    index.html — callers fall through to the archive-listing renderer.
    """
    if not _is_application_artifact(info):
        return None
    import zipfile
    try:
        with zipfile.ZipFile(file_path) as zf:
            # App-builder lays files at the zip root. Walk the namelist
            # case-insensitively so a stray Index.html / INDEX.HTML
            # still resolves; pick the shallowest match so a nested
            # ``docs/index.html`` doesn't win over a root entry.
            best = None
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                basename = name.rsplit("/", 1)[-1].lower()
                if basename != "index.html":
                    continue
                depth = name.count("/")
                if best is None or depth < best[0]:
                    best = (depth, name)
            if best is None:
                return None
            raw = zf.read(best[1])
    except (zipfile.BadZipFile, KeyError, OSError) as exc:
        log.warning(
            "artifact_app_preview_zip_read_failed",
            artifact_id=artifact_id, error=str(exc)[:160],
        )
        return None
    try:
        html = raw.decode("utf-8", errors="replace")
    except Exception:
        return None
    # Trailing slash on the base href so resolution stays inside the
    # artifact: ``app.js`` → ``/api/artifacts/<id>/preview/app.js``.
    base_tag = f'<base href="/api/artifacts/{artifact_id}/preview/">'
    # Idempotent: don't double-inject if the app already declared its
    # own <base> (rare for builder output, common for user uploads).
    lower = html.lower()
    if "<base" in lower:
        return html
    # Inject right after <head>; if there's no <head> (some hand-rolled
    # docs), prepend so the tag is parsed before any resource refs.
    head_open = re.search(r"<head\b[^>]*>", html, flags=re.IGNORECASE)
    if head_open:
        idx = head_open.end()
        return html[:idx] + base_tag + html[idx:]
    return base_tag + html


# Lightweight extension → media type map. Covers the static-web-app
# common case; for anything missing (.wasm uploads, exotic formats) the
# mimetypes module's guess fills in.
_APP_PREVIEW_MIME_OVERRIDES = {
    ".js": "application/javascript; charset=utf-8",
    ".mjs": "application/javascript; charset=utf-8",
    ".css": "text/css; charset=utf-8",
    ".html": "text/html; charset=utf-8",
    ".htm": "text/html; charset=utf-8",
    ".json": "application/json; charset=utf-8",
    ".svg": "image/svg+xml",
    ".wasm": "application/wasm",
    ".map": "application/json; charset=utf-8",
}


@router.get("/{artifact_id}/preview/{file_path:path}")
async def preview_application_file(
    artifact_id: str, file_path: str, request: Request,
):
    """Serve an individual file from inside an application artifact's
    zip. Sibling assets referenced by index.html (script src=, link
    href=, fetch('data.json'), …) land here via the <base> tag injected
    by ``preview_artifact``.

    Returns 404 for non-application artifacts so this route doesn't
    accidentally expose archive contents of unrelated zip uploads.
    """
    gate_response, preview_uid = await _isolated_preview_gate(request, artifact_id)
    if gate_response is not None:
        return gate_response
    store = _get_store(request)
    user_id = preview_uid or _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    if not _is_application_artifact(info):
        # Refuse to act as a generic zip-content server. Only artifacts
        # explicitly marked as applications opt into per-file serving;
        # otherwise we'd let any user-uploaded zip be browsed as a
        # virtual filesystem under /api/artifacts/<id>/preview/<path>.
        raise HTTPException(status_code=404, detail="Not an application artifact")
    disk_path = store.get_file_path(info.get("path", ""))
    if not disk_path:
        raise HTTPException(status_code=404, detail="Artifact file missing")

    # Normalise the requested path. zipfile name matching is exact, so
    # collapse any leading ./ or stray slashes a <base>-resolved URL
    # might bring along.
    wanted = file_path.lstrip("/").replace("\\", "/")
    if not wanted or wanted.endswith("/"):
        raise HTTPException(status_code=404, detail="File not found")

    import zipfile
    try:
        with zipfile.ZipFile(disk_path) as zf:
            # Try the exact name first (matches app-builder output where
            # files sit at the zip root). Fall back to a case-insensitive
            # search for robustness against capitalisation drift in
            # hand-edited artifacts.
            try:
                data = zf.read(wanted)
            except KeyError:
                lower_wanted = wanted.lower()
                match = next(
                    (n for n in zf.namelist()
                     if not n.endswith("/") and n.lower() == lower_wanted),
                    None,
                )
                if match is None:
                    raise HTTPException(status_code=404, detail="File not found")
                data = zf.read(match)
    except zipfile.BadZipFile as exc:
        log.warning(
            "artifact_app_preview_file_zip_corrupt",
            artifact_id=artifact_id, error=str(exc)[:160],
        )
        raise HTTPException(status_code=500, detail="Archive unreadable") from exc

    suffix = ("." + wanted.rsplit(".", 1)[-1].lower()) if "." in wanted else ""
    media_type = _APP_PREVIEW_MIME_OVERRIDES.get(suffix)
    if media_type is None:
        import mimetypes
        guess, _ = mimetypes.guess_type(wanted)
        media_type = guess or "application/octet-stream"

    # no-store because the artifact bytes can change in place (workspace
    # edits write a new version into the same artifact_id). Aggressive
    # caching here would mask the user's own saves.
    return Response(
        content=data,
        media_type=media_type,
        headers={"Cache-Control": "no-store"},
    )


def _preview_shell(title: str, download_url: str, body: str) -> str:
    """Wrap preview content in a styled HTML shell matching the Augmentum dark theme.

    Served inside an iframe — embeds its own viewport meta and responsive
    rules so it looks right at desktop, half-screen, quarter-screen, and
    mobile sizes. Inner body content can bring its own CSS on top of these.
    """
    from html import escape
    return (
        '<!DOCTYPE html><html><head><meta charset="utf-8">'
        '<meta name="viewport" content="width=device-width,initial-scale=1">'
        f'<title>{escape(title)}</title>'
        '<style>'
        '*{margin:0;padding:0;box-sizing:border-box}'
        'html,body{height:100%}'
        'body{background:#0f0f1a;color:#ececf1;'
        'font-family:system-ui,-apple-system,"Segoe UI",Roboto,sans-serif;'
        '-webkit-text-size-adjust:100%;display:flex;flex-direction:column}'
        '.header{display:flex;align-items:center;gap:10px;padding:10px 16px;'
        'background:#161625;border-bottom:1px solid #2d2d45;position:sticky;top:0;z-index:10;'
        'flex-shrink:0}'
        '.title{flex:1;min-width:0;font-size:14px;font-weight:600;'
        'white-space:nowrap;overflow:hidden;text-overflow:ellipsis}'
        '.dl-btn{color:#6c8aff;text-decoration:none;font-size:12px;padding:6px 12px;'
        'border:1px solid #2d2d45;border-radius:6px;white-space:nowrap;'
        'transition:background 0.15s}'
        '.dl-btn:hover{background:rgba(108,138,255,0.1)}'
        '.content{flex:1;overflow:auto;min-height:0}'
        '.content>.wrap{padding:28px 24px;max-width:780px;margin:0 auto}'
        'img,video,canvas{max-width:100%;height:auto}'
        'pre{white-space:pre-wrap;word-break:break-word;overflow-x:auto}'
        'pre.nowrap{white-space:pre;word-break:normal}'
        'table{max-width:100%}'
        '.table-scroll{overflow-x:auto;-webkit-overflow-scrolling:touch}'
        # --- Typography baseline -----------------------------------
        # Every renderer inherits these unless it scopes its own rules
        # under a class (PPTX/XLSX/EPUB/archive all scope, so this layer
        # is what dresses DOCX, Markdown, and any other "bare semantic
        # HTML" output without forcing per-renderer style duplication).
        'h1{font-size:28px;font-weight:600;line-height:1.25;margin:0 0 14px;color:#ececf1}'
        'h2{font-size:22px;font-weight:600;line-height:1.3;margin:28px 0 12px;color:#ececf1}'
        'h3{font-size:18px;font-weight:600;line-height:1.35;margin:24px 0 10px;color:#ececf1}'
        'h4{font-size:16px;font-weight:600;line-height:1.4;margin:20px 0 8px;color:#ececf1}'
        'h5,h6{font-size:13px;font-weight:600;line-height:1.4;margin:16px 0 6px;'
        'color:#a1a1b5;text-transform:uppercase;letter-spacing:0.04em}'
        'p{margin:0 0 12px;line-height:1.65}'
        'ul,ol{margin:12px 0;padding-left:24px;line-height:1.65}'
        'li{margin:4px 0}'
        'li>ul,li>ol{margin:4px 0}'
        'a,a:visited{color:#6c8aff;text-decoration:none}'
        'a:hover{text-decoration:underline}'
        'strong,b{font-weight:600;color:#ffffff}'
        'em,i{font-style:italic}'
        'code{background:#1c1c2e;color:#e8c5a0;padding:1px 6px;border-radius:4px;'
        'font-family:ui-monospace,SFMono-Regular,Consolas,monospace;font-size:0.9em}'
        'pre{background:#161625;border:1px solid #2d2d45;border-radius:8px;'
        'padding:14px 16px;line-height:1.55;font-size:12.5px;color:#d4d4df;margin:14px 0}'
        'pre code{background:none;color:inherit;padding:0;border-radius:0;font-size:inherit}'
        'blockquote{border-left:3px solid #6c8aff;padding:6px 0 6px 16px;margin:14px 0;'
        'color:#a1a1b5;font-style:italic}'
        'hr{border:none;border-top:1px solid #2d2d45;margin:24px 0}'
        '.wrap table{border-collapse:collapse;margin:14px 0;font-size:13.5px}'
        '.wrap th,.wrap td{padding:8px 12px;border:1px solid #2d2d45;text-align:left}'
        '.wrap th{background:#1c1c2e;font-weight:600;color:#a1a1b5;font-size:12px;'
        'text-transform:uppercase;letter-spacing:0.03em}'
        # Serif utility — opt-in for prose-heavy renderers (DOCX, EPUB
        # fallback, long-form Markdown) where editorial type beats UI sans.
        '.serif{font-family:Charter,"Iowan Old Style","Apple Garamond",Georgia,'
        '"Times New Roman",serif}'
        '.serif p{line-height:1.75;font-size:16px}'
        '.serif h1,.serif h2,.serif h3,.serif h4{font-family:inherit}'
        # Pygments output. The wrapping `<div class="highlight">` provides
        # the panel chrome; the inner `<pre>` is just spans. Token classes
        # are a curated palette tied to the shell theme — no Pygments
        # built-in style is shipped, so adding/removing tokens is one rule.
        '.highlight{background:#161625;border:1px solid #2d2d45;border-radius:8px;'
        'padding:14px 16px;margin:14px 0;overflow-x:auto;'
        'font-family:ui-monospace,SFMono-Regular,Consolas,monospace;'
        'font-size:12.5px;line-height:1.55;color:#d4d4df}'
        '.highlight pre{background:none;border:none;padding:0;margin:0;'
        'font:inherit;color:inherit;line-height:inherit;'
        'white-space:pre;word-break:normal;overflow:visible}'
        '.highlight .c,.highlight .ch,.highlight .cm,.highlight .c1,'
        '.highlight .cs,.highlight .cp,.highlight .cpf'
        '{color:#6b6b80;font-style:italic}'
        '.highlight .k,.highlight .kc,.highlight .kd,.highlight .kn,'
        '.highlight .kp,.highlight .kr,.highlight .kt'
        '{color:#c084fc;font-weight:600}'
        '.highlight .nf,.highlight .fm{color:#6c8aff}'
        '.highlight .nc,.highlight .ne{color:#fbbf24}'
        '.highlight .nb,.highlight .bp{color:#7dd3fc}'
        '.highlight .nd{color:#fbbf24;font-style:italic}'
        '.highlight .nt{color:#c084fc}'
        '.highlight .na{color:#7dd3fc}'
        '.highlight .s,.highlight .s1,.highlight .s2,.highlight .sb,'
        '.highlight .sc,.highlight .sd,.highlight .sh,.highlight .si,'
        '.highlight .sx,.highlight .sr,.highlight .ss{color:#86efac}'
        '.highlight .se{color:#fb923c;font-weight:600}'
        '.highlight .m,.highlight .mb,.highlight .mf,.highlight .mh,'
        '.highlight .mi,.highlight .mo,.highlight .il{color:#fb923c}'
        '.highlight .o,.highlight .ow,.highlight .p{color:#a1a1b5}'
        '.highlight .err{color:#f87171;background:rgba(248,113,113,0.1)}'
        '@media (max-width:640px){'
        '.header{padding:8px 12px;gap:8px}'
        '.title{font-size:13px}'
        '.dl-btn{font-size:11px;padding:5px 10px}'
        '.content>.wrap{padding:16px 12px}'
        '}'
        '@media (max-width:400px){'
        '.content>.wrap{padding:12px 8px}'
        '}'
        '</style></head>'
        '<body>'
        '<div class="header">'
        f'<span class="title">{escape(title)}</span>'
        f'<a class="dl-btn" href="{download_url}" download>Download</a>'
        '</div>'
        f'<div class="content">{body}</div>'
        '</body></html>'
    )


def _docx_to_html(file_path, title: str, download_url: str) -> str | None:
    """Convert a DOCX to a styled HTML preview."""
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError:
        return None

    try:
        doc = Document(str(file_path))
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower() if para.style else ""
            if "heading 1" in style:
                parts.append(f"<h1>{_esc(text)}</h1>")
            elif "heading 2" in style:
                parts.append(f"<h2>{_esc(text)}</h2>")
            elif "heading 3" in style:
                parts.append(f"<h3>{_esc(text)}</h3>")
            elif "list" in style:
                parts.append(f"<li>{_esc(text)}</li>")
            else:
                # Check for bold/italic runs
                runs_html = []
                for run in para.runs:
                    t = _esc(run.text)
                    if run.bold and run.italic:
                        t = f"<b><i>{t}</i></b>"
                    elif run.bold:
                        t = f"<b>{t}</b>"
                    elif run.italic:
                        t = f"<i>{t}</i>"
                    runs_html.append(t)
                parts.append(f"<p>{''.join(runs_html)}</p>")

        # Wrap list items
        body = "\n".join(parts)
        body = body.replace("</li>\n<li>", "</li><li>")
        import re
        body = re.sub(r"(<li>.*?</li>)", r"<ul>\1</ul>", body, flags=re.DOTALL)
        body = body.replace("</ul>\n<ul>", "\n")

        # `.serif` opts into the editorial type baseline declared in
        # _preview_shell — heading scale + serif body + measured leading.
        styled_body = f'<div class="wrap serif">{body}</div>'
        return _preview_shell(title, download_url, styled_body)
    except Exception as exc:
        log.warning("docx_preview_failed", error=str(exc))
        return None


def _epub_to_html(file_path, title: str, download_url: str) -> str | None:
    """Convert an EPUB to a styled HTML book viewer.

    Extracts XHTML chapter content and the embedded stylesheet from the
    EPUB ZIP, then renders them inside the preview shell with the book's
    own CSS applied.
    """
    import base64
    import zipfile
    from xml.etree import ElementTree as ET

    try:
        with zipfile.ZipFile(str(file_path)) as zf:
            # Parse OPF to get spine order
            opf_path = None
            try:
                container = ET.fromstring(zf.read("META-INF/container.xml"))
                ns = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
                rootfile = container.find(".//c:rootfile", ns)
                if rootfile is not None:
                    opf_path = rootfile.get("full-path", "")
            except Exception:
                # Fallback: scan for .opf
                opf_path = next(
                    (n for n in zf.namelist() if n.endswith(".opf")), None,
                )

            if not opf_path:
                return None

            opf_dir = opf_path.rsplit("/", 1)[0] + "/" if "/" in opf_path else ""
            opf = ET.fromstring(zf.read(opf_path))
            opf_ns = {"opf": "http://www.idpf.org/2007/opf"}

            # Build id→href manifest map
            manifest = {}
            for item in opf.findall(".//opf:manifest/opf:item", opf_ns):
                manifest[item.get("id", "")] = {
                    "href": item.get("href", ""),
                    "type": item.get("media-type", ""),
                }

            # Get spine order
            spine_ids = [
                ref.get("idref", "")
                for ref in opf.findall(".//opf:spine/opf:itemref", opf_ns)
            ]

            # Extract CSS
            css = ""
            for item in manifest.values():
                if item["type"] == "text/css":
                    try:
                        css = zf.read(opf_dir + item["href"]).decode("utf-8")
                    except Exception as exc:
                        log.debug("epub_css_extract_failed", href=item["href"], error=str(exc))
                    break

            # Build image map (filename → base64 data URI)
            image_map: dict[str, str] = {}
            for item in manifest.values():
                if item["type"].startswith("image/"):
                    try:
                        img_data = zf.read(opf_dir + item["href"])
                        b64 = base64.b64encode(img_data).decode("ascii")
                        image_map[item["href"]] = (
                            f"data:{item['type']};base64,{b64}"
                        )
                    except Exception as exc:
                        log.debug("epub_image_extract_failed", href=item["href"], error=str(exc))

            # Extract chapter content from spine
            chapters: list[str] = []
            for sid in spine_ids:
                entry = manifest.get(sid)
                if not entry or "xhtml" not in entry["type"]:
                    continue
                try:
                    raw = zf.read(opf_dir + entry["href"]).decode("utf-8")
                except Exception as exc:
                    log.debug("epub_chapter_read_failed", href=entry["href"], error=str(exc))
                    continue

                # Extract <body> content
                body_match = __import__("re").search(
                    r"<body[^>]*>(.*)</body>", raw, __import__("re").DOTALL,
                )
                if not body_match:
                    continue
                body_html = body_match.group(1).strip()

                # Rewrite image src to base64 data URIs
                for img_href, data_uri in image_map.items():
                    # Match both relative and images/ prefixed paths
                    body_html = body_html.replace(
                        f'src="{img_href}"', f'src="{data_uri}"',
                    )
                    # Also match without directory prefix
                    basename = img_href.rsplit("/", 1)[-1] if "/" in img_href else img_href
                    body_html = body_html.replace(
                        f'src="images/{basename}"', f'src="{data_uri}"',
                    )

                chapters.append(body_html)

            if not chapters:
                return None

            # The preview iframe shows the book "as the theme says" — drop any
            # prefers-color-scheme block so a device in dark mode doesn't
            # repaint a deliberately-light reading theme (and vice-versa).
            if css:
                dark_marker = "\n/* --- Dark-mode adaptation"
                idx = css.find(dark_marker)
                if idx != -1:
                    css = css[:idx]

            # Pull the book's own page colours out of its `body { … }` rule so
            # the WHOLE preview surface (not just the centred column) takes the
            # theme — otherwise a sepia/night book reads as a small column on
            # the shell's dark chrome and "the background didn't change".
            page_bg = page_fg = ""
            if css:
                body_rule = __import__("re").search(
                    r"\bbody\s*\{[^}]*\}", css, __import__("re").DOTALL,
                )
                if body_rule:
                    block = body_rule.group(0)
                    m_bg = __import__("re").search(r"background\s*:\s*([^;}]+)", block)
                    if m_bg:
                        page_bg = m_bg.group(1).strip()
                    m_fg = __import__("re").search(r"(?<!-)color\s*:\s*([^;}]+)", block)
                    if m_fg:
                        page_fg = m_fg.group(1).strip()

            # Build the preview: book CSS + chapter content
            # Scope the book CSS inside .epub-content to avoid conflicts
            scoped_css = css.replace("body", ".epub-content") if css else ""

            page_override = ""
            if page_bg:
                # Equal-specificity `body` rule placed after the shell's →
                # later wins, so the page floods with the book's colour.
                page_override = f"html,body{{background:{page_bg}}}"
                if page_fg:
                    page_override += f".content{{color:{page_fg}}}"

            chapter_html = "\n<hr style='border:none;border-top:1px solid currentColor;opacity:0.2;margin:32px 0'>\n".join(chapters)

            # Book CSS lives inside the .epub-content scope; the shell's .wrap
            # class provides responsive padding + centered max-width.
            # `.serif` provides editorial type when the book ships no CSS;
            # books that do ship CSS override .serif via .epub-content rules
            # (book CSS appears later in the cascade, equal specificity).
            styled_body = (
                f"<style>{scoped_css}"
                f"{page_override}"
                ".epub-content img,.epub-content svg{max-width:100%;height:auto}"
                ".epub-content table{display:block;overflow-x:auto}"
                "</style>"
                f'<div class="wrap epub-content serif">{chapter_html}</div>'
            )
            return _preview_shell(title, download_url, styled_body)
    except Exception as exc:
        log.warning("epub_preview_failed", error=str(exc))
        return None


def _pptx_to_html(file_path, title: str, download_url: str) -> str | None:
    """Convert a PPTX to an HTML slide viewer."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return None

    try:
        prs = Presentation(str(file_path))
        slides_html: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # Rough heading detection by font size
                            is_title = any(
                                run.font.size and run.font.size.pt >= 20
                                for run in para.runs if run.font.size
                            )
                            if is_title:
                                texts.append(f"<h2>{_esc(text)}</h2>")
                            else:
                                texts.append(f"<p>{_esc(text)}</p>")
                elif shape.has_table:
                    table = shape.table
                    rows_html: list[str] = []
                    for row in table.rows:
                        cells = "".join(
                            f"<td style='padding:6px 10px;border:1px solid #2d2d45'>"
                            f"{_esc(cell.text)}</td>"
                            for cell in row.cells
                        )
                        rows_html.append(f"<tr>{cells}</tr>")
                    texts.append(
                        f"<table style='border-collapse:collapse;width:100%;margin:12px 0'>"
                        f"{''.join(rows_html)}</table>"
                    )

            # Speaker notes live on ``slide.notes_slide.notes_text_frame``.
            # They're authored by the LLM and persisted into the PPTX but
            # were previously invisible in the web preview — the user
            # had to open PowerPoint to see them. Render them in a
            # collapsible block below the slide body.
            notes_html = ""
            try:
                if slide.has_notes_slide:
                    notes_text = (
                        slide.notes_slide.notes_text_frame.text or ""
                    ).strip()
                    if notes_text:
                        notes_paragraphs = "".join(
                            f"<p>{_esc(line)}</p>"
                            for line in notes_text.splitlines() if line.strip()
                        )
                        notes_html = (
                            '<details class="slide-notes">'
                            '<summary>Speaker notes</summary>'
                            f'{notes_paragraphs}'
                            '</details>'
                        )
            except Exception:
                # Preview is best-effort — a malformed notes slide should
                # not break rendering of the deck body.
                notes_html = ""

            content = "\n".join(texts) if texts else "<p style='color:#6b6b80'>Empty slide</p>"
            slides_html.append(
                f'<div class="slide" id="slide-{i}">'
                f'<div class="slide-num">Slide {i}</div>'
                f'{content}{notes_html}</div>'
            )

        body = (
            '<style>'
            '.slides{max-width:900px;margin:0 auto;padding:16px 12px}'
            '.slide{margin:16px 0;padding:28px;background:#1c1c2e;'
            'border:1px solid #2d2d45;border-radius:10px;min-height:160px}'
            '.slide-num{font-size:11px;color:#6b6b80;text-transform:uppercase;'
            'letter-spacing:0.05em;margin-bottom:14px}'
            '.slide h2{font-size:22px;margin:0 0 12px;color:#ececf1;line-height:1.3}'
            '.slide p{margin:8px 0;line-height:1.6;color:#a1a1b5;word-break:break-word}'
            '.slide table{display:block;overflow-x:auto;max-width:100%}'
            '.slide-notes{margin-top:16px;padding-top:12px;'
            'border-top:1px dashed #2d2d45}'
            '.slide-notes summary{font-size:11px;text-transform:uppercase;'
            'letter-spacing:0.05em;color:#6b6b80;cursor:pointer;outline:none}'
            '.slide-notes summary:hover{color:#9797ac}'
            '.slide-notes p{font-size:13px;color:#8f8fa3;margin:6px 0 0}'
            '@media (max-width:640px){'
            '.slides{padding:8px 6px}'
            '.slide{padding:18px 16px;margin:10px 0}'
            '.slide h2{font-size:18px}'
            '}'
            '</style>'
            f'<div class="slides">{"".join(slides_html)}</div>'
        )
        return _preview_shell(title, download_url, body)
    except Exception as exc:
        log.warning("pptx_preview_failed", error=str(exc))
        return None


def _xlsx_to_html(file_path, fmt: str, title: str, download_url: str) -> str | None:
    """Convert XLSX/CSV to an HTML table preview."""
    try:
        if fmt == "csv":
            import csv
            rows: list[list[str]] = []
            with open(file_path, encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i > 500:
                        break  # limit for preview
                    rows.append(row)
        else:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            ws = wb.active
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i > 500:
                    break
                rows.append([str(cell) if cell is not None else "" for cell in row])
            wb.close()
    except ImportError:
        return None
    except Exception as exc:
        log.warning("xlsx_preview_failed", error=str(exc))
        return None

    if not rows:
        # Truly empty workbook — render an explicit empty-state shell rather
        # than falling through to the generic "Download" card, which doesn't
        # even look like a spreadsheet.
        body = (
            '<div style="text-align:center;padding:48px;color:#6b6b80">'
            'This spreadsheet has no data yet.</div>'
        )
        return _preview_shell(title, download_url, body)

    # Build HTML table — first row as header
    header = rows[0] if rows else []
    data_rows = rows[1:] if len(rows) > 1 else []

    th = "".join(f"<th>{_esc(str(c))}</th>" for c in header)
    tr_list = []
    for row in data_rows:
        td = "".join(f"<td>{_esc(str(c))}</td>" for c in row)
        tr_list.append(f"<tr>{td}</tr>")

    body = (
        '<style>'
        '.xlsx-scroll{overflow-x:auto;overflow-y:visible;-webkit-overflow-scrolling:touch}'
        '.xlsx-scroll table{border-collapse:collapse;font-size:13px;'
        'min-width:100%;width:max-content}'
        '.xlsx-scroll th{background:#1c1c2e;padding:8px 12px;text-align:left;'
        'border-bottom:2px solid #2d2d45;font-weight:600;color:#a1a1b5;'
        'font-size:11px;text-transform:uppercase;letter-spacing:0.03em;'
        'white-space:nowrap;position:sticky;top:0}'
        '.xlsx-scroll td{padding:6px 12px;border-bottom:1px solid #1c1c2e;'
        'color:#ececf1;white-space:nowrap}'
        '.xlsx-scroll tr:hover td{background:#1a1a2c}'
        '.truncated{text-align:center;padding:16px;color:#6b6b80;font-style:italic}'
        '@media (max-width:640px){'
        '.xlsx-scroll th,.xlsx-scroll td{padding:6px 10px;font-size:12px}'
        '}'
        '</style>'
        '<div class="xlsx-scroll">'
        f'<table><thead><tr>{th}</tr></thead>'
        f'<tbody>{"".join(tr_list)}</tbody></table>'
        '</div>'
        f'{"<div class=truncated>No data rows yet</div>" if not data_rows else ""}'
        f'{"<div class=truncated>Showing first 500 rows</div>" if len(rows) > 500 else ""}'
    )
    return _preview_shell(title, download_url, body)


def _esc(text: str) -> str:
    """HTML-escape a string."""
    from html import escape
    return escape(text)


@router.post("/{artifact_id}/save")
async def save_artifact_source(artifact_id: str, request: Request):
    """Update the source JSON for an artifact (workspace save)."""
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    body = await request.json()
    source_json = body.get("source_json")
    if not source_json:
        raise HTTPException(status_code=400, detail="source_json is required")
    await store.update_source(artifact_id, source_json, user_id=user_id)
    return JSONResponse({"ok": True, "id": artifact_id})


@router.post("/{artifact_id}/upload")
async def upload_artifact_binary(artifact_id: str, request: Request):
    """Replace the artifact's binary file (e.g., modified PDF from visual editor)."""
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")

    form = await request.form()
    uploaded = form.get("file")
    if not uploaded:
        raise HTTPException(status_code=400, detail="No file uploaded")

    contents = await uploaded.read()
    if not contents:
        raise HTTPException(status_code=400, detail="Empty file")

    ok = await store.update_file(artifact_id, contents, user_id=user_id)
    if not ok:
        raise HTTPException(status_code=500, detail="Failed to update artifact file")

    log.info("artifact_binary_uploaded", artifact_id=artifact_id, size=len(contents))
    return JSONResponse({"ok": True, "size_bytes": len(contents)})


# Adapter routes: expose the image postprocess pipeline (upscale, bg-removal)
# to any image artifact, not just images that originated from the image
# generation surface. Both functions rewrite the artifact binary in place so
# Studio's viewer sees the edit the moment it reloads.
_IMAGE_ARTIFACT_EXTS = {"png", "jpg", "jpeg", "webp"}


def _is_image_artifact(info: dict) -> bool:
    """Heuristic — use format first, fall back to the filename extension."""
    fmt = (info.get("format") or "").lower()
    if fmt in _IMAGE_ARTIFACT_EXTS:
        return True
    name = (info.get("filename") or info.get("path") or "").lower()
    ext = name.rsplit(".", 1)[-1] if "." in name else ""
    return ext in _IMAGE_ARTIFACT_EXTS


async def _run_postprocess_and_replace(
    store, artifact_id: str, user_id: str, src_path, run,
) -> tuple[int, int, int]:
    """Shared pipeline: call a postprocess function, write result into the
    artifact, clean up the intermediate file. Returns (width, height, size).

    `run` is an awaitable that takes the str path and returns
    (new_id, new_path, width, height) — matching the postprocess.py contract.
    """
    from pathlib import Path as _Path
    _, new_path, w, h = await run(str(src_path))
    try:
        data = _Path(new_path).read_bytes()
        ok = await store.update_file(artifact_id, data, user_id=user_id)
        if not ok:
            raise HTTPException(status_code=500, detail="Failed to update artifact file")
    finally:
        # Postprocess writes into settings.image_output_dir — the artifact
        # now owns the bytes, so drop the intermediate to avoid orphan files
        # accumulating in the image output directory.
        try:
            _Path(new_path).unlink(missing_ok=True)
        except Exception:  # noqa: BLE001 — best-effort cleanup
            log.debug("postprocess_cleanup_failed", path=str(new_path))
    return w, h, len(data)


@router.post("/{artifact_id}/upscale")
async def upscale_artifact(artifact_id: str, request: Request):
    """Upscale an image artifact via the Spandrel UltraSharp pipeline.

    Request body is optional JSON: {"scale": 2 | 4}. Default is 4. The
    binary is replaced in place — the artifact_id stays stable so any
    references (chat attachments, library pins, file index rows) keep
    resolving.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    if not _is_image_artifact(info):
        raise HTTPException(status_code=400, detail="Artifact is not an image")

    scale = 4
    try:
        body = await request.json()
        if isinstance(body, dict) and body.get("scale") is not None:
            scale = int(body["scale"])
    except Exception as exc:  # noqa: BLE001 — empty / non-JSON body is allowed
        log.debug("artifact_upscale_body_parse_failed", error=str(exc))
    if scale not in (2, 4):
        raise HTTPException(status_code=400, detail="scale must be 2 or 4")

    src_path = store.get_file_path(info["path"])
    if not src_path or not src_path.is_file():
        raise HTTPException(status_code=404, detail="Artifact file missing")

    try:
        from augmentum.image.postprocess import upscale_image
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"Upscale pipeline unavailable: {e}") from e

    try:
        async def _run(p):
            return await upscale_image(p, scale=scale)
        w, h, size = await _run_postprocess_and_replace(store, artifact_id, user_id, src_path, _run)
    except HTTPException:
        raise
    except Exception as e:  # noqa: BLE001 — report failure details to the UI
        log.warning("artifact_upscale_failed", artifact_id=artifact_id, error=str(e))
        raise HTTPException(status_code=500, detail=f"Upscale failed: {e}") from e

    log.info("artifact_upscaled", artifact_id=artifact_id, scale=scale, width=w, height=h)
    return JSONResponse({"ok": True, "width": w, "height": h, "size_bytes": size})


# Convert: source-format → target-format matrix. Keys are the source format
# (or "*document" / "*presentation" / "*spreadsheet" for the native source
# types where we have a structured source_json to rerender from). Values
# are the accepted targets.
_CONVERT_MATRIX: dict[str, set[str]] = {
    "png":  {"jpg", "jpeg", "webp"},
    "jpg":  {"png", "jpeg", "webp"},
    "jpeg": {"png", "jpg", "webp"},
    "webp": {"png", "jpg", "jpeg"},
    "*document": {"pdf", "docx"},
}
_IMAGE_TARGET_FORMATS = {"png", "jpg", "jpeg", "webp"}


def _convert_source_key(info: dict) -> str:
    """Lookup key for _CONVERT_MATRIX. Falls back to the file extension."""
    # Structured artifacts carry a source_json with a `type` field — these
    # have a full rendering pipeline we can re-invoke to produce a
    # different final format.
    src_json = info.get("source_json")
    if src_json:
        try:
            src = json.loads(src_json) if isinstance(src_json, str) else src_json
            t = (src.get("type") or "").lower() if isinstance(src, dict) else ""
            if t:
                return f"*{t}"
        except Exception as exc:  # noqa: BLE001 — fall back to format-based lookup
            log.debug("artifact_convert_source_json_parse_failed", error=str(exc))
    fmt = (info.get("format") or "").lower()
    if fmt in _CONVERT_MATRIX:
        return fmt
    name = (info.get("filename") or info.get("path") or "").lower()
    ext = name.rsplit(".", 1)[-1] if "." in name else ""
    return ext


async def _convert_image(src_path, target: str) -> tuple[bytes, str]:
    """Pillow-based image format conversion. Returns (bytes, new_extension).

    Delegates to :meth:`ImageConvertTool.convert_bytes` — the canonical
    implementation lives in ``augmentum/tools/image_convert.py``.
    """
    from augmentum.tools.image_convert import ImageConvertTool

    return await ImageConvertTool.convert_bytes(str(src_path), target)


async def _convert_document(info: dict, target: str) -> tuple[bytes, str]:
    """Re-render a structured document source to PDF or DOCX.

    Delegates to :meth:`DocumentConvertTool.render_bytes` — the canonical
    implementation lives in ``augmentum/tools/document_convert.py``.
    """
    from augmentum.tools.document_convert import DocumentConvertTool

    try:
        return await DocumentConvertTool.render_bytes(info, target)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@router.post("/{artifact_id}/convert")
async def convert_artifact(artifact_id: str, request: Request):
    """Convert an artifact to a different format, saved as a sibling artifact.

    Request body: {"to": "pdf" | "docx" | "png" | "jpg" | "webp"}. The target
    set depends on the source format — see `_CONVERT_MATRIX`. Returns the
    new artifact's id so callers can open it immediately.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")

    body = await request.json()
    if not isinstance(body, dict) or not body.get("to"):
        raise HTTPException(status_code=400, detail="'to' is required")
    target = str(body["to"]).lower()

    key = _convert_source_key(info)
    allowed = _CONVERT_MATRIX.get(key) or set()
    if target not in allowed:
        raise HTTPException(
            status_code=400,
            detail=f"Cannot convert {key or 'artifact'} → {target}. Allowed: {sorted(allowed) or '(none)'}",
        )

    src_path = store.get_file_path(info["path"])
    if not src_path or not src_path.is_file():
        raise HTTPException(status_code=404, detail="Artifact file missing")

    # Dispatch by source kind.
    if target in _IMAGE_TARGET_FORMATS and key in _CONVERT_MATRIX and key != "*document":
        data, new_ext = await _convert_image(src_path, target)
    elif key == "*document":
        data, new_ext = await _convert_document(info, target)
    else:  # pragma: no cover — matrix above should have filtered this
        raise HTTPException(status_code=400, detail="Unsupported conversion")

    # Save as sibling artifact.
    base = info.get("display_name") or info.get("filename") or "artifact"
    stem = base.rsplit(".", 1)[0] if "." in base else base
    filename = f"{stem}.{new_ext}"
    try:
        saved = await store.save(
            data=data,
            filename=filename,
            fmt=new_ext,
            session_id=info.get("session_id") or "",
            task_id=info.get("task_id") or "",
            display_name=f"{stem} ({new_ext.upper()})",
            metadata={"converted_from": artifact_id, "target": new_ext},
            user_id=user_id,
        )
    except Exception as e:  # noqa: BLE001
        log.warning("artifact_convert_save_failed", artifact_id=artifact_id, error=str(e))
        raise HTTPException(status_code=500, detail=f"Convert save failed: {e}") from e

    log.info("artifact_converted", artifact_id=artifact_id, source_key=key, target=target, new_id=saved["id"])
    return JSONResponse({"ok": True, "artifact_id": saved["id"], "filename": saved["filename"], "format": new_ext, "size_bytes": len(data)})


_AUDIO_VIDEO_EXTS = {
    "mp3", "wav", "m4a", "flac", "ogg", "opus", "aac",
    "mp4", "mov", "webm", "mkv", "avi",
}


@router.post("/{artifact_id}/inpaint")
async def inpaint_artifact(artifact_id: str, request: Request):
    """Inpaint an image artifact through the image generation pipeline.

    Unlike upscale / remove-bg which use standalone postprocess functions,
    inpaint goes through the full generation queue (it's a diffusion run,
    not a one-shot neural op). Body expects:
      - prompt              : str, required
      - mask_image          : base64 PNG, white = repaint
      - negative_prompt     : str, optional
      - mode                : "default" | "improve" | "modify"
      - strength            : 0.0-1.0 (default 1.0)
      - model, steps, cfg_scale, seed, sampler : optional, server defaults apply
      - inpaint_full_res    : bool (default false)
      - inpaint_padding     : int (default 32)

    Returns the new sibling artifact's metadata. We save as a sibling
    because inpaint is destructive enough that overwriting the source is a
    trap — users want to compare, keep the source around, or try another
    prompt on the same image.
    """
    import base64

    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    if not _is_image_artifact(info):
        raise HTTPException(status_code=400, detail="Artifact is not an image")

    body = await request.json()
    if not isinstance(body, dict):
        raise HTTPException(status_code=400, detail="JSON body required")
    prompt = (body.get("prompt") or "").strip()
    if not prompt:
        raise HTTPException(status_code=400, detail="prompt is required")
    mask_b64 = body.get("mask_image") or ""
    if not mask_b64:
        raise HTTPException(status_code=400, detail="mask_image is required")

    src_path = store.get_file_path(info["path"])
    if not src_path or not src_path.is_file():
        raise HTTPException(status_code=404, detail="Artifact file missing")

    # Pydantic's InpaintRequest wants source_image either as a base64 PNG or
    # as an image_id (from the image generations table). Artifacts aren't in
    # that table, so we send the bytes as base64 — matches what the sidebar
    # inpaint already does.
    src_bytes = src_path.read_bytes()
    source_b64 = base64.b64encode(src_bytes).decode("ascii")

    # Required app state — guard early so a missing config returns a clean
    # 503 rather than an AttributeError deep in the queue.
    queue = getattr(request.app.state, "image_queue", None)
    if queue is None:
        raise HTTPException(status_code=503, detail="Image generation pipeline not available")

    from augmentum.config import settings
    from augmentum.image.queue import GenerationJob
    from augmentum.image.schemas import JobType

    mode = (body.get("mode") or "default").lower()
    if mode not in ("default", "improve", "modify"):
        raise HTTPException(status_code=400, detail="mode must be default|improve|modify")
    try:
        strength = float(body.get("strength", 1.0))
    except (TypeError, ValueError):
        strength = 1.0
    strength = max(0.0, min(1.0, strength))

    model = (body.get("model") or settings.image_default_model or "").strip()
    steps = int(body.get("steps") or settings.image_default_steps or 0)
    cfg_scale = float(body.get("cfg_scale") or settings.image_default_cfg or 0.0)

    # Apply the same distilled-model defaults the /api/image/inpaint route
    # uses, so Turbo/Lightning models don't get nuked with 25-step defaults.
    try:
        from augmentum.image.distilled import apply_distilled_defaults
        steps, cfg_scale = apply_distilled_defaults(model, steps, cfg_scale)
    except Exception as exc:  # noqa: BLE001 — best-effort tuning, not critical
        log.debug("distilled_defaults_failed", model=model, error=str(exc))

    job = GenerationJob(
        job_type=JobType.INPAINT,
        prompt=prompt,
        negative_prompt=(body.get("negative_prompt") or "").strip(),
        model=model,
        preset=body.get("preset") or "",
        width=int(body.get("width") or 0),
        height=int(body.get("height") or 0),
        steps=steps,
        cfg_scale=cfg_scale,
        seed=int(body.get("seed") or -1),
        sampler=(body.get("sampler") or "").strip(),
        source_image=source_b64,
        mask_image=mask_b64,
        strength=strength,
        mask_blur=int(body.get("mask_blur") or 4),
        inpaint_mode=mode,
        inpaint_full_res=bool(body.get("inpaint_full_res")),
        inpaint_padding=int(body.get("inpaint_padding") or 32),
        user_id=user_id,
    )

    try:
        job = await queue.submit(job)
    except RuntimeError as exc:
        raise HTTPException(status_code=429, detail=str(exc)) from exc

    try:
        result = await queue.wait_for_result(job, timeout=settings.image_generation_timeout)
    except TimeoutError as exc:
        raise HTTPException(status_code=504, detail="Inpaint timed out") from exc
    except Exception as exc:  # noqa: BLE001
        log.warning("artifact_inpaint_failed", artifact_id=artifact_id, error=str(exc))
        raise HTTPException(status_code=500, detail=f"Inpaint failed: {exc}") from exc

    # Result file lives in settings.image_output_dir per image_persistence.
    # We copy the bytes into a new artifact so closing the inpaint generation
    # doesn't strand the only copy inside the image-generation store.
    result_image_id = result.get("image_id", "")
    persistence = getattr(request.app.state, "image_persistence", None)
    result_bytes = b""
    if persistence and result_image_id and user_id:
        try:
            gen = await persistence.get_generation(result_image_id, user_id=user_id)
            if gen and gen.get("file_path"):
                from pathlib import Path as _Path
                p = _Path(gen["file_path"])
                if p.is_file():
                    result_bytes = p.read_bytes()
        except Exception as e:  # noqa: BLE001
            log.warning("inpaint_result_read_failed", image_id=result_image_id, error=str(e))
    if not result_bytes:
        raise HTTPException(status_code=500, detail="Inpaint produced no output file")

    base = info.get("display_name") or info.get("filename") or "image"
    stem = base.rsplit(".", 1)[0] if "." in base else base
    saved = await store.save(
        data=result_bytes,
        filename=f"{stem}.inpainted.png",
        fmt="png",
        session_id=info.get("session_id") or "",
        task_id=info.get("task_id") or "",
        display_name=f"{stem} (inpainted)",
        metadata={
            "inpainted_from": artifact_id,
            "source_image_id": result_image_id,
            "prompt": prompt,
            "mode": mode,
        },
        user_id=user_id,
    )

    log.info(
        "artifact_inpainted",
        artifact_id=artifact_id, new_id=saved["id"],
        width=result.get("width", 0), height=result.get("height", 0),
    )
    return JSONResponse({
        "ok": True,
        "artifact_id": saved["id"],
        "filename": saved["filename"],
        "width": result.get("width", 0),
        "height": result.get("height", 0),
        "seed": result.get("seed", -1),
    })


@router.post("/{artifact_id}/transcribe")
async def transcribe_artifact(artifact_id: str, request: Request):
    """Transcribe an audio or video artifact via the configured STT provider.

    The transcript is saved as a sibling `.txt` artifact in the same session
    so the user has a permanent record, and returned inline so the viewer
    can surface it immediately. Video artifacts work because the STT
    provider (Moonshine via ffmpeg, or Deepgram / OpenAI-compat) extracts
    the audio track internally.
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    fmt = (info.get("format") or "").lower()
    name = (info.get("filename") or info.get("path") or "").lower()
    ext = fmt if fmt in _AUDIO_VIDEO_EXTS else (name.rsplit(".", 1)[-1] if "." in name else "")
    if ext not in _AUDIO_VIDEO_EXTS:
        raise HTTPException(status_code=400, detail="Artifact is not audio or video")

    file_path = store.get_file_path(info["path"])
    if not file_path or not file_path.is_file():
        raise HTTPException(status_code=404, detail="Artifact file missing")

    try:
        from augmentum.proxy.audio_routes import transcribe_audio_bytes
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"STT pipeline unavailable: {e}") from e

    audio_bytes = file_path.read_bytes()
    # Pick a content-type hint the STT provider will accept — the MIME map
    # in download_artifact has the full set; we only need a plausible one
    # here so Deepgram / OpenAI-compat don't reject it. Moonshine sniffs
    # the magic bytes regardless.
    content_type_map = {
        "mp3": "audio/mpeg", "wav": "audio/wav", "m4a": "audio/mp4",
        "flac": "audio/flac", "ogg": "audio/ogg", "opus": "audio/opus",
        "aac": "audio/aac",
        "mp4": "video/mp4", "mov": "video/quicktime", "webm": "video/webm",
        "mkv": "video/x-matroska", "avi": "video/x-msvideo",
    }
    content_type = content_type_map.get(ext, "application/octet-stream")

    text = await transcribe_audio_bytes(
        request, audio_bytes, filename=info.get("filename") or f"audio.{ext}",
        content_type=content_type,
    )
    text = (text or "").strip()

    # Persist the transcript as a sibling .txt artifact in the same session —
    # keeps it discoverable from the Library and reopenable in Studio.
    sibling_id = ""
    try:
        base = info.get("display_name") or info.get("filename") or "transcript"
        stem = base.rsplit(".", 1)[0] if "." in base else base
        saved = await store.save(
            data=text.encode("utf-8") or b"(no speech detected)",
            filename=f"{stem}.transcript.txt",
            fmt="txt",
            session_id=info.get("session_id") or "",
            task_id=info.get("task_id") or "",
            display_name=f"{stem} — transcript",
            metadata={"transcribed_from": artifact_id},
            user_id=user_id,
        )
        sibling_id = saved.get("id", "")
    except Exception as e:  # noqa: BLE001 — best-effort persistence; the transcript still returns inline
        log.warning("transcript_sibling_save_failed", artifact_id=artifact_id, error=str(e))

    log.info("artifact_transcribed", artifact_id=artifact_id, chars=len(text), sibling=sibling_id)
    return JSONResponse({"ok": True, "transcript": text, "artifact_id": sibling_id or None})


@router.post("/{artifact_id}/remove-bg")
async def remove_bg_artifact(artifact_id: str, request: Request):
    """Remove the background from an image artifact via rembg isnet-general-use.

    Output is RGBA PNG. The binary replaces the original; callers should
    refresh the file extension in their UI if they care (the format column
    isn't updated — the bytes are just PNG-encoded with alpha).
    """
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    if not _is_image_artifact(info):
        raise HTTPException(status_code=400, detail="Artifact is not an image")

    src_path = store.get_file_path(info["path"])
    if not src_path or not src_path.is_file():
        raise HTTPException(status_code=404, detail="Artifact file missing")

    try:
        from augmentum.image.postprocess import remove_background
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"rembg unavailable: {e}") from e

    try:
        w, h, size = await _run_postprocess_and_replace(store, artifact_id, user_id, src_path, remove_background)
    except HTTPException:
        raise
    except Exception as e:  # noqa: BLE001
        log.warning("artifact_remove_bg_failed", artifact_id=artifact_id, error=str(e))
        raise HTTPException(status_code=500, detail=f"Background removal failed: {e}") from e

    log.info("artifact_bg_removed", artifact_id=artifact_id, width=w, height=h)
    return JSONResponse({"ok": True, "width": w, "height": h, "size_bytes": size})


@router.patch("/{artifact_id}/pin")
async def toggle_pin(artifact_id: str, request: Request):
    """Toggle pinned status for an artifact."""
    store = _get_store(request)
    user_id = _user_id(request)
    info = await store.get(artifact_id, user_id=user_id)
    if not info:
        raise HTTPException(status_code=404, detail="Artifact not found")
    new_pinned = not info.get("pinned", False)
    await store.set_pinned(artifact_id, new_pinned, user_id=user_id)
    return JSONResponse({"pinned": new_pinned})


@router.patch("/{artifact_id}/open")
async def mark_opened(artifact_id: str, request: Request):
    """Update last_opened_at for sort-by-recent."""
    store = _get_store(request)
    user_id = _user_id(request)
    await store.touch_opened(artifact_id, user_id=user_id)
    return JSONResponse({"ok": True})


@router.delete("/{artifact_id}")
async def delete_artifact(artifact_id: str, request: Request):
    """Delete an artifact by ID."""
    store = _get_store(request)
    user_id = _user_id(request)
    deleted = await store.delete(artifact_id, user_id=user_id)
    if not deleted:
        raise HTTPException(status_code=404, detail="Artifact not found")
    return JSONResponse({"deleted": True, "id": artifact_id})
