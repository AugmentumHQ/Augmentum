"""Presentation artifact tool — generates PPTX files from structured slides."""

from __future__ import annotations

import io
import json
import os
import re
from typing import TYPE_CHECKING

from augmentum.tools.base import Tool, ToolCategory, ToolResult
from augmentum.utils.logging import get_logger

if TYPE_CHECKING:
    from augmentum.tools.artifact_storage import ArtifactStore

log = get_logger(__name__)


class PresentationTool(Tool):
    """Generate PowerPoint presentations from structured slide definitions."""

    def __init__(self, artifact_store: ArtifactStore) -> None:
        self._store = artifact_store

    @property
    def name(self) -> str:
        return "create_presentation"

    @property
    def description(self) -> str:
        return (
            "Create a PowerPoint presentation (.pptx) from structured slides. "
            "Each slide has a layout type (title, content, two-column, or blank), "
            "a title, body content, and optional speaker notes. "
            "Returns a download link for the generated file."
            " Write slide content directly — no meta-commentary or placeholder text."
        )

    @property
    def category(self) -> ToolCategory:
        return ToolCategory.ARTIFACT

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Presentation title (used on title slide and metadata)",
                },
                "subtitle": {
                    "type": "string",
                    "description": "Subtitle for the title slide (optional)",
                    "default": "",
                },
                "author": {
                    "type": "string",
                    "description": "Author name (optional)",
                    "default": "",
                },
                "theme": {
                    "type": "string",
                    "description": "Visual theme preset (slate, corporate, modern, emerald, rose). Default: slate",
                    "default": "",
                },
                "slides": {
                    "type": "array",
                    "description": "Ordered list of slides",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {
                                "type": "string",
                                "enum": ["title", "content", "two_column", "blank"],
                                "description": "Slide layout type (default: content)",
                                "default": "content",
                            },
                            "title": {
                                "type": "string",
                                "description": "Slide title",
                            },
                            "body": {
                                "type": "string",
                                "description": (
                                    "Slide body text. Use bullet lines (- item) "
                                    "for bullet points. For two_column layout, "
                                    "separate columns with '|||'."
                                ),
                            },
                            "notes": {
                                "type": "string",
                                "description": "Speaker notes for this slide (optional)",
                                "default": "",
                            },
                            "image_url": {
                                "type": "string",
                                "description": (
                                    "URL of an image to embed on this slide "
                                    "(e.g. /api/image/abc123 from image_generation tool)"
                                ),
                                "default": "",
                            },
                            "additional_images": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": (
                                    "Up to 3 extra image URLs (post-render picker "
                                    "append). Rendered as bottom-right thumbnails."
                                ),
                                "default": [],
                            },
                        },
                        "required": ["title"],
                    },
                },
            },
            "required": ["title", "slides"],
        }

    async def execute(
        self,
        *,
        title: str = "Presentation",
        subtitle: str = "",
        author: str = "",
        theme: str = "",
        slides: list | None = None,
        task_id: str = "",
        session_id: str = "",
        **kwargs,
    ) -> ToolResult:
        from augmentum.tools.artifact_normalize import normalize_slides, normalize_str

        title = normalize_str(title, "Presentation")
        subtitle = normalize_str(subtitle)
        author = normalize_str(author)
        slides = normalize_slides(slides)
        if not slides:
            return ToolResult(success=False, error="No slides provided")

        from augmentum.tools.artifact_sanitize import sanitize_slides
        slides = sanitize_slides(slides)
        if not slides:
            return ToolResult(success=False, error="No slides after cleanup")

        # Resolve image paths from URLs before rendering
        resolved_slides = await self._resolve_images(slides)

        try:
            data = _render_pptx(title, subtitle, author, resolved_slides, theme_name=theme)

            safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')[:60]
            filename = f"{safe_title}.pptx"

            image_count = sum(
                1 for s in resolved_slides if s.get("_image_path")
            )

            info = await self._store.save(
                data=data,
                filename=filename,
                fmt="pptx",
                task_id=task_id,
                session_id=session_id,
                display_name=f"{title}.pptx",
                user_id=Tool.extract_user_id(kwargs),
                metadata={
                    "page_type": "presentation",
                    "slide_count": len(slides),
                    "image_count": image_count,
                    "author": author,
                },
                source_json=json.dumps({
                    "type": "presentation",
                    "title": title,
                    "subtitle": subtitle,
                    "author": author,
                    "theme": theme,
                    "slides": slides,
                }),
            )

            from augmentum.tools.base import make_artifact_card

            img_str = f", {image_count} image{'s' if image_count != 1 else ''}" if image_count else ""
            summary = (
                f"Presentation '{title}' is ready — "
                f"{len(slides)} slide{'s' if len(slides) != 1 else ''}{img_str}. "
                "Available in the artifact library."
            )
            card = make_artifact_card(
                info,
                kind="artifact",
                title=title,
                subtitle=subtitle or (f"by {author}" if author else "PPTX"),
                summary=summary,
                preview={
                    "artifact_kind": "presentation",
                    "format": "pptx",
                    "size_bytes": info.get("size_bytes", 0),
                    "slides": [
                        {"heading": s.get("heading", f"Slide {i+1}")}
                        for i, s in enumerate(slides)
                    ],
                    "image_count": image_count,
                },
            )
            return ToolResult(
                success=True,
                output=summary,
                metadata=info,
                card=card,
            )
        except Exception as e:
            log.error("presentation_creation_failed", error=str(e), exc_info=True)
            return ToolResult(success=False, error=f"Presentation creation failed: {e}")


    async def _resolve_images(self, slides: list) -> list:
        """Resolve image_url and additional_images references to filesystem paths."""
        resolved = []
        for slide in slides:
            s = dict(slide)
            image_url = s.get("image_url", "")
            if image_url:
                path = await self._resolve_image_path(image_url)
                if path:
                    s["_image_path"] = str(path)
                else:
                    log.warning("image_resolve_failed", url=image_url)
            # additional_images = bottom-right thumbnails. Resolve each;
            # drop any that fail rather than logging once per slide.
            extras = s.get("additional_images") or []
            extra_paths: list[str] = []
            for url in extras[:3]:  # picker hard caps at 3
                if not url:
                    continue
                p = await self._resolve_image_path(url)
                if p:
                    extra_paths.append(str(p))
            if extra_paths:
                s["_additional_image_paths"] = extra_paths
            resolved.append(s)
        return resolved

    async def _resolve_image_path(self, url: str) -> str | None:
        """Resolve an image URL to a local filesystem path."""
        m = re.match(r"/api/image/([a-zA-Z0-9_-]+)", url)
        if m:
            image_id = m.group(1)
            return await self._resolve_image_id(image_id)

        m = re.match(r"/api/artifacts/([a-zA-Z0-9_-]+)/download", url)
        if m:
            artifact_id = m.group(1)
            info = await self._store.get(artifact_id)
            if info:
                file_path = self._store.get_file_path(info["path"])
                if file_path:
                    return str(file_path)
        return None

    async def _resolve_image_id(self, image_id: str) -> str | None:
        """Look up an image_id via the image_generations table."""
        try:
            db = self._store._db
            async with db.execute(
                "SELECT file_path FROM image_generations WHERE image_id = ?",
                (image_id,),
            ) as cursor:
                row = await cursor.fetchone()
                if row and row[0] and os.path.exists(row[0]):
                    return row[0]
        except Exception as e:
            log.debug("image_id_lookup_failed", image_id=image_id, error=str(e))
        return None


# ---------------------------------------------------------------------------
# PPTX rendering via python-pptx
# ---------------------------------------------------------------------------

# Slide layout indices in the default template
_LAYOUT_TITLE = 0       # Title Slide
_LAYOUT_CONTENT = 1     # Title and Content
_LAYOUT_TWO_COL = 3     # Two Content
_LAYOUT_BLANK = 6       # Blank

def _get_pptx_theme(theme_name: str = "", design: dict | None = None):
    """Get theme colors for PPTX rendering.

    `design` (optional) is the Studio design block; when present, font sizes
    + accent are first applied via `apply_design` so the returned dict has
    the user-customised values already baked in. Renderer paths can stay
    as-is — they just see different numbers when the user picked them.
    """
    from augmentum.tools.artifact_theme import FONT_FAMILY_PPTX, apply_design, get_theme
    theme = apply_design(get_theme(theme_name), design)
    family_pref = (design or {}).get("font_family", "system")
    return {
        "accent": theme.rgb(theme.accent),
        "accent_dark": theme.rgb(theme.accent_dark),
        "text": theme.rgb(theme.text),
        "text_secondary": theme.rgb(theme.text_secondary),
        "text_muted": theme.rgb(theme.text_muted),
        "border": theme.rgb(theme.border),
        "bar_height": theme.slide_bar_height,
        "title_size": theme.slide_title_size,
        "body_size": theme.slide_body_size,
        "font_name": FONT_FAMILY_PPTX.get(family_pref, "Calibri"),
    }


def _add_slide_bar(slide, prs, colors):
    """Add a thin accent bar at the bottom of a content slide."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    bar_h = colors["bar_height"]
    if bar_h <= 0:
        return

    left = Inches(0)
    top = prs.slide_height - Inches(bar_h)
    width = prs.slide_width
    height = Inches(bar_h)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*colors["accent"])
    shape.line.fill.background()  # no border


def _render_pptx(
    title: str, subtitle: str, author: str, slides: list,
    theme_name: str = "",
    design: dict | None = None,
) -> bytes:
    """Render slides to a PPTX file using python-pptx + theme system.

    `design` applies font_family + font_size_scale + accent_override before
    rendering. Renderer code is unchanged — `colors` already encoded font
    size + accent and now carries `font_name` too.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    colors = _get_pptx_theme(theme_name, design)
    font_name = colors.get("font_name", "Calibri")

    prs = Presentation()

    # Set slide dimensions (widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Core properties
    prs.core_properties.title = title
    if author:
        prs.core_properties.author = author

    # -- Title slide --
    title_layout = prs.slide_layouts[_LAYOUT_TITLE]
    slide = prs.slides.add_slide(title_layout)

    # Title placeholder
    if slide.placeholders[0]:
        tf = slide.placeholders[0].text_frame
        tf.text = title
        for para in tf.paragraphs:
            para.font.size = Pt(40)
            para.font.bold = True
            para.font.color.rgb = RGBColor(*colors["text"])

    # Subtitle placeholder
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        tf = slide.placeholders[1].text_frame
        parts = []
        if subtitle:
            parts.append(subtitle)
        if author:
            parts.append(author)
        parts.append("Generated by Augmentum")
        tf.text = "\n".join(parts)
        for para in tf.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(*colors["text_muted"])

    # -- Content slides --
    for slide_def in slides:
        layout_name = slide_def.get("layout", "content")
        slide_title = slide_def.get("title", "")
        body = slide_def.get("body", "")
        notes = slide_def.get("notes", "")
        image_path = slide_def.get("_image_path", "")
        extra_paths = slide_def.get("_additional_image_paths") or []

        if layout_name == "title":
            _add_section_title_slide(prs, slide_title, body, colors)
        elif layout_name == "two_column":
            slide = _add_two_column_slide(prs, slide_title, body, notes, image_path, colors)
            _embed_additional_thumbnails(slide, prs, extra_paths)
        elif layout_name == "blank":
            slide = _add_blank_slide(prs, slide_title, body, notes, image_path, colors)
            _embed_additional_thumbnails(slide, prs, extra_paths)
        else:
            slide = _add_content_slide(prs, slide_title, body, notes, image_path, colors)
            _embed_additional_thumbnails(slide, prs, extra_paths)

    # design.font_family — apply once at the end so every text frame's
    # paragraphs pick up the chosen face. python-pptx propagates font.name
    # from paragraph → run when runs don't override it, so this is enough
    # without rewriting every per-paragraph block above.
    if font_name and font_name != "Calibri":
        _apply_font_name_to_presentation(prs, font_name)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _apply_font_name_to_presentation(prs, font_name: str) -> None:
    """Set font.name on every text frame's paragraphs in the presentation.

    Walks slides → shapes → text_frames → paragraphs. Idempotent — runs
    inherit from paragraphs in python-pptx so we don't need to recurse into
    individual runs.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                if para.font.name is None:
                    para.font.name = font_name


def _add_content_slide(prs, title: str, body: str, notes: str, image_path: str = "",
                       colors: dict | None = None):
    """Add a standard title + content slide. Returns the created slide."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    c = colors or _get_pptx_theme()

    layout = prs.slide_layouts[_LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)

    # Title
    if slide.placeholders[0]:
        tf = slide.placeholders[0].text_frame
        tf.text = title
        for para in tf.paragraphs:
            para.font.size = Pt(c["title_size"])
            para.font.bold = True
            para.font.color.rgb = RGBColor(*c["text"])

    # Body content
    if body and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        _populate_text_frame(tf, body)

    # Embed image
    if image_path and os.path.exists(image_path):
        _embed_slide_image(slide, prs, image_path, has_body=bool(body))

    # Accent bar at bottom
    _add_slide_bar(slide, prs, c)

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _add_section_title_slide(prs, title: str, body: str, colors: dict | None = None) -> None:
    """Add a section divider slide (uses title layout)."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    c = colors or _get_pptx_theme()

    layout = prs.slide_layouts[_LAYOUT_TITLE]
    slide = prs.slides.add_slide(layout)

    if slide.placeholders[0]:
        tf = slide.placeholders[0].text_frame
        tf.text = title
        for para in tf.paragraphs:
            para.font.size = Pt(36)
            para.font.bold = True
            para.font.color.rgb = RGBColor(*c["accent"])

    if body and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = body
        for para in tf.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(*c["text_muted"])

    _add_slide_bar(slide, prs, c)


def _add_two_column_slide(prs, title: str, body: str, notes: str, image_path: str = "",
                          colors: dict | None = None):
    """Add a two-column slide. Returns the created slide."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    c = colors or _get_pptx_theme()

    columns = body.split("|||", 1) if body else ["", ""]
    left = columns[0].strip()
    right = columns[1].strip() if len(columns) > 1 else ""

    layout = prs.slide_layouts[_LAYOUT_TWO_COL]
    slide = prs.slides.add_slide(layout)

    if slide.placeholders[0]:
        tf = slide.placeholders[0].text_frame
        tf.text = title
        for para in tf.paragraphs:
            para.font.size = Pt(c["title_size"])
            para.font.bold = True
            para.font.color.rgb = RGBColor(*c["text"])

    # Left column (placeholder index 1)
    if left and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        _populate_text_frame(tf, left)

    # Right column (placeholder index 2)
    if right and len(slide.placeholders) > 2:
        tf = slide.placeholders[2].text_frame
        tf.clear()
        _populate_text_frame(tf, right)

    # Embed image
    if image_path and os.path.exists(image_path):
        _embed_slide_image(slide, prs, image_path, has_body=bool(body))

    _add_slide_bar(slide, prs, c)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _add_blank_slide(prs, title: str, body: str, notes: str, image_path: str = "",
                     colors: dict | None = None):
    """Add a blank slide. Returns the created slide."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    c = colors or _get_pptx_theme()

    layout = prs.slide_layouts[_LAYOUT_BLANK]
    slide = prs.slides.add_slide(layout)

    content = title
    if body:
        content = f"{title}\n\n{body}" if title else body

    if content:
        left = Inches(1)
        top = Inches(1.5)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(3)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # First line as title
        lines = content.split("\n", 1)
        para = tf.paragraphs[0]
        para.text = lines[0]
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = RGBColor(*c["text"])
        para.alignment = PP_ALIGN.CENTER

        if len(lines) > 1:
            body_para = tf.add_paragraph()
            body_para.text = lines[1].strip()
            body_para.font.size = Pt(16)
            body_para.font.color.rgb = RGBColor(*c["text_muted"])
            body_para.alignment = PP_ALIGN.CENTER

    # Embed image on blank slide
    if image_path and os.path.exists(image_path):
        _embed_slide_image(slide, prs, image_path, has_body=bool(content))

    _add_slide_bar(slide, prs, c)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# ---------------------------------------------------------------------------
# Image embedding helper
# ---------------------------------------------------------------------------

def _embed_slide_image(slide, prs, image_path: str, *, has_body: bool = False) -> None:
    """Embed an image on a slide, positioned based on whether body text exists."""
    from pptx.util import Inches

    try:
        from PIL import Image as PILImage

        with PILImage.open(image_path) as img:
            img_w, img_h = img.size

        slide_w = prs.slide_width.inches
        slide_h = prs.slide_height.inches

        if has_body:
            # Place image on the right side, vertically centered
            max_w = slide_w * 0.4
            max_h = slide_h * 0.6
            scale = min(max_w / img_w, max_h / img_h, 1.0)
            display_w = img_w * scale
            display_h = img_h * scale
            left = Inches(slide_w - display_w - 0.5)
            top = Inches((slide_h - display_h) / 2)
        else:
            # No body text — center the image large
            max_w = slide_w * 0.8
            max_h = slide_h * 0.7
            scale = min(max_w / img_w, max_h / img_h, 1.0)
            display_w = img_w * scale
            display_h = img_h * scale
            left = Inches((slide_w - display_w) / 2)
            top = Inches((slide_h - display_h) / 2 + 0.3)

        slide.shapes.add_picture(
            image_path,
            left, top,
            Inches(display_w), Inches(display_h),
        )
    except Exception as e:
        log.warning("pptx_image_embed_failed", path=image_path, error=str(e))


# ---------------------------------------------------------------------------
# Additional-image thumbnail strip (post-render picker "append" support)
# ---------------------------------------------------------------------------

def _embed_additional_thumbnails(slide, prs, paths: list[str]) -> None:
    """Lay out additional_images as a bottom-right thumbnail strip.

    Spec (Phase 1): primary image stays in its current slot; appended
    images render as 1.0×1.0 inch thumbnails right-aligned along a row
    above the accent bar. Max 3 — the picker enforces the same cap.
    Body collision is accepted for v1; real composition is Phase 4.
    """
    if not paths:
        return
    from pptx.util import Inches

    keep = [p for p in paths[:3] if p and os.path.exists(p)]
    if not keep:
        return
    thumb_w = 1.0
    thumb_h = 1.0
    gap = 0.1
    accent_bar_inset = 0.2  # leave the bottom accent bar visible
    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches
    strip_w = len(keep) * thumb_w + (len(keep) - 1) * gap
    left_start = slide_w - strip_w - 0.3
    top = slide_h - thumb_h - accent_bar_inset - 0.1
    for i, path in enumerate(keep):
        try:
            slide.shapes.add_picture(
                path,
                Inches(left_start + i * (thumb_w + gap)),
                Inches(top),
                Inches(thumb_w),
                Inches(thumb_h),
            )
        except Exception as exc:
            log.warning("pptx_additional_thumbnail_failed",
                        path=path, error=str(exc))


# ---------------------------------------------------------------------------
# Text frame population (shared helper)
# ---------------------------------------------------------------------------

def _populate_text_frame(tf, body: str) -> None:
    """Populate a text frame with body text, handling bullets and paragraphs."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    lines = body.strip().split("\n")
    first = True

    for line in lines:
        stripped = line.strip()
        if not stripped:
            # Empty line — add spacing paragraph
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            para.text = ""
            first = False
            continue

        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        # Bullet detection
        is_bullet = False
        indent_level = 0
        text = stripped

        if stripped.startswith(("- ", "* ", "• ")):
            is_bullet = True
            text = stripped[2:].strip()
        elif re.match(r'^\d+[.)]\s', stripped):
            is_bullet = True
            text = re.sub(r'^\d+[.)]\s*', '', stripped)
        elif stripped.startswith("  - ") or stripped.startswith("  * "):
            is_bullet = True
            indent_level = 1
            text = stripped.strip()[2:].strip()

        # Apply bold/italic from markdown
        _apply_formatted_runs(para, text)

        para.font.size = Pt(16)
        # Use near-black from the default theme for text frame content
        para.font.color.rgb = RGBColor(15, 23, 42)  # Slate-900

        if is_bullet:
            para.level = indent_level
            # python-pptx uses level to auto-apply bullet styling
            # in content placeholders


def _apply_formatted_runs(paragraph, text: str) -> None:
    """Parse basic markdown bold/italic and add formatted runs to a paragraph."""

    # Split on **bold** and *italic* markers
    parts = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith("*") and part.endswith("*"):
            run = paragraph.add_run()
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run = paragraph.add_run()
            run.text = part
