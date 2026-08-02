"""Document parsing tool — extracts text from PDF, DOCX, PPTX, XLSX files."""

from __future__ import annotations

import asyncio
import os

from augmentum.tools.base import Tool, ToolCategory, ToolResult
from augmentum.utils.logging import get_logger

log = get_logger(__name__)

# Maximum characters to return from parsed content.
_MAX_OUTPUT_CHARS = 50_000

# File extensions to parser function mapping.
_SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".csv", ".md"}


def _parse_pdf(path: str) -> str:
    """Extract text from a PDF using pdfplumber."""
    try:
        import pdfplumber
    except ImportError:
        raise ImportError(
            "pdfplumber is not installed. Install with: pip install pdfplumber"
        )

    pages: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                pages.append(f"--- Page {i} ---\n{text}")

            # Extract tables if present.
            tables = page.extract_tables()
            for table in tables:
                if table:
                    rows = []
                    for row in table:
                        cells = [str(c) if c is not None else "" for c in row]
                        rows.append(" | ".join(cells))
                    pages.append(f"[Table on page {i}]\n" + "\n".join(rows))

    return "\n\n".join(pages) if pages else "(No extractable text in PDF)"


def _parse_docx(path: str) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs) if paragraphs else "(No text content in DOCX)"


def _parse_pptx(path: str) -> str:
    """Extract text from a PPTX presentation."""
    from pptx import Presentation

    prs = Presentation(path)
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    return "\n\n".join(slides) if slides else "(No text content in presentation)"


def _parse_xlsx(path: str) -> str:
    """Extract data from an XLSX spreadsheet."""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    sheets: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()

    return "\n\n".join(sheets) if sheets else "(No data in spreadsheet)"


def _parse_text(path: str) -> str:
    """Read a plain text file."""
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".txt": _parse_text,
    ".csv": _parse_text,
    ".md": _parse_text,
}


class DocumentParseTool(Tool):
    """Parse and extract text from documents (PDF, DOCX, PPTX, XLSX, TXT)."""

    @property
    def name(self) -> str:
        return "document_parse"

    @property
    def description(self) -> str:
        return (
            "Parse a document file and extract its text content. "
            "Supports PDF (with tables), DOCX, PPTX, XLSX, TXT, CSV, and Markdown files."
        )

    @property
    def category(self) -> ToolCategory:
        return ToolCategory.FILE

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the document file to parse",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Maximum characters to return from the parsed content",
                    "default": 50000,
                },
            },
            "required": ["path"],
        }

    @property
    def timeout(self) -> float:
        return 30.0

    @property
    def cacheable(self) -> bool:
        return True

    @property
    def cache_ttl(self) -> float:
        return 600.0  # 10 minutes

    def __init__(self, base_dir: str | None = None) -> None:
        self._base_dir = base_dir

    def _resolve_path(self, path: str) -> str:
        """Resolve and validate the file path, preventing path traversal."""
        if self._base_dir:
            # Resolve relative to base_dir, prevent traversal.
            resolved = os.path.normpath(os.path.join(self._base_dir, path))
            base = os.path.normpath(self._base_dir)
            if not resolved.startswith(base + os.sep) and resolved != base:
                raise ValueError(f"Path traversal blocked: {path}")
            return resolved
        return os.path.abspath(path)

    def validate_input(self, **kwargs) -> bool:
        path = kwargs.get("path", "")
        return isinstance(path, str) and len(path.strip()) > 0

    async def execute(
        self,
        *,
        path: str,
        max_chars: int = 50000,
    ) -> ToolResult:
        """Parse a document and extract its text content."""
        if not path.strip():
            return ToolResult(success=False, error="Empty file path")

        try:
            resolved = self._resolve_path(path)
        except ValueError as exc:
            return ToolResult(success=False, error=str(exc))

        if not os.path.exists(resolved):
            return ToolResult(
                success=False,
                error=f"File not found: {path}",
            )

        ext = os.path.splitext(resolved)[1].lower()
        parser = _PARSERS.get(ext)
        if parser is None:
            return ToolResult(
                success=False,
                error=f"Unsupported file type '{ext}'. "
                      f"Supported: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}",
            )

        try:
            content = await asyncio.to_thread(parser, resolved)
        except ImportError as exc:
            return ToolResult(success=False, error=str(exc))
        except Exception as exc:
            log.warning(
                "document_parse_failed",
                path=path,
                ext=ext,
                error=str(exc),
            )
            return ToolResult(
                success=False,
                error=f"Failed to parse {ext} file: {exc}",
            )

        if not content:
            return ToolResult(
                success=True,
                output="(Document contains no extractable text)",
                metadata={"path": path, "type": ext, "char_count": 0},
            )

        max_chars = min(max_chars, _MAX_OUTPUT_CHARS)
        truncated = len(content) > max_chars
        if truncated:
            content = content[:max_chars] + "\n\n... (truncated)"

        return ToolResult(
            success=True,
            output=content,
            metadata={
                "path": path,
                "type": ext,
                "char_count": len(content),
                "truncated": truncated,
            },
        )
